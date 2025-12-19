"""
Document indexing module.

Handles extraction and indexing of text from various document formats.
Supports PDF, DOCX, XLSX, PPTX, HTML, text files, and unknown formats via hex extraction.
"""

import logging
import concurrent.futures
from typing import Optional, Callable, Dict, Any, List, Tuple
from pathlib import Path
from datetime import datetime
import traceback

from .db import DatabaseManager
from .utils import (
    calculate_sha256,
    detect_file_type,
    is_text_file,
    should_ignore,
    format_size,
)
from .hex_extractor import extract_text_from_unknown

logger = logging.getLogger(__name__)


class TextExtractor:
    """Extract text from various document formats."""

    def __init__(self, trust_external_tools: bool = False):
        """
        Initialize text extractor.

        Args:
            trust_external_tools: Whether to use external conversion tools
        """
        self.trust_external_tools = trust_external_tools

    def extract(self, file_path: Path, file_type: str) -> Tuple[str, str]:
        """
        Extract text from a file.

        Args:
            file_path: Path to file
            file_type: File type/extension

        Returns:
            Tuple of (extracted_text, extractor_name)
        """
        # Map file types to extraction methods
        extractors = {
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "doc": self._extract_doc,
            "xlsx": self._extract_xlsx,
            "xls": self._extract_xls,
            "pptx": self._extract_pptx,
            "ppt": self._extract_ppt,
            "html": self._extract_html,
            "htm": self._extract_html,
            "xml": self._extract_xml,
            "txt": self._extract_text,
            "md": self._extract_text,
            "rst": self._extract_text,
            "csv": self._extract_text,
            "json": self._extract_text,
            "log": self._extract_text,
            "py": self._extract_text,
            "js": self._extract_text,
            "java": self._extract_text,
            "c": self._extract_text,
            "cpp": self._extract_text,
            "h": self._extract_text,
            "cs": self._extract_text,
            "go": self._extract_text,
            "rs": self._extract_text,
            "rb": self._extract_text,
            "php": self._extract_text,
            "sh": self._extract_text,
            "bat": self._extract_text,
            "ps1": self._extract_text,
        }

        extractor = extractors.get(file_type, self._extract_unknown)
        return extractor(file_path)

    def _extract_pdf(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from PDF."""
        try:
            from pdfminer.high_level import extract_text

            text = extract_text(str(file_path))
            return text, "pdfminer"
        except Exception as e:
            logger.error(f"PDF extraction failed for {file_path}: {e}")
            # Fallback to hex extraction
            return self._extract_unknown(file_path)

    def _extract_docx(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from DOCX."""
        try:
            from docx import Document

            doc = Document(str(file_path))
            text = "\n".join([para.text for para in doc.paragraphs])
            return text, "python-docx"
        except Exception as e:
            logger.error(f"DOCX extraction failed for {file_path}: {e}")
            return self._extract_unknown(file_path)

    def _extract_doc(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from DOC (old Word format)."""
        # Old .doc format requires external tools
        logger.warning(
            f"Old .doc format not supported, using hex extraction: {file_path}"
        )
        return self._extract_unknown(file_path)

    def _extract_xlsx(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from XLSX."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)

            text_parts = []
            for sheet in wb.worksheets:
                text_parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(
                        [str(cell) if cell is not None else "" for cell in row]
                    )
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts), "openpyxl"

        except Exception as e:
            logger.error(f"XLSX extraction failed for {file_path}: {e}")
            return self._extract_unknown(file_path)

    def _extract_xls(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from XLS (old Excel format)."""
        logger.warning(
            f"Old .xls format not supported, using hex extraction: {file_path}"
        )
        return self._extract_unknown(file_path)

    def _extract_pptx(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[Slide {slide_num}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_parts.append(shape.text)

            return "\n".join(text_parts), "python-pptx"

        except Exception as e:
            logger.error(f"PPTX extraction failed for {file_path}: {e}")
            return self._extract_unknown(file_path)

    def _extract_ppt(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from PPT (old PowerPoint format)."""
        logger.warning(
            f"Old .ppt format not supported, using hex extraction: {file_path}"
        )
        return self._extract_unknown(file_path)

    def _extract_html(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from HTML."""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "rb") as f:
                soup = BeautifulSoup(f, "lxml")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator="\n", strip=True)
            return text, "beautifulsoup4"

        except Exception as e:
            logger.error(f"HTML extraction failed for {file_path}: {e}")
            return self._extract_text(file_path)

    def _extract_xml(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from XML."""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "rb") as f:
                soup = BeautifulSoup(f, "lxml-xml")

            text = soup.get_text(separator="\n", strip=True)
            return text, "beautifulsoup4"

        except Exception as e:
            logger.error(f"XML extraction failed for {file_path}: {e}")
            return self._extract_text(file_path)

    def _extract_text(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from plain text file."""
        try:
            # Try multiple encodings
            for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        text = f.read()
                    return text, f"text ({encoding})"
                except UnicodeDecodeError:
                    continue

            # If all encodings fail, use binary mode with error handling
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            return text, "text (utf-8 with errors)"

        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            return "", "error"

    def _extract_unknown(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from unknown format using hex extraction."""
        try:
            text = extract_text_from_unknown(file_path)
            return text, "hex-extractor"
        except Exception as e:
            logger.error(f"Hex extraction failed for {file_path}: {e}")
            return "", "error"


class IndexProgress:
    """Progress tracking for indexing operations."""

    def __init__(self):
        self.total_files = 0
        self.processed_files = 0
        self.successful = 0
        self.failed = 0
        self.skipped = 0
        self.bytes_processed = 0
        self.current_file = ""
        self.start_time = datetime.now()
        self.errors: List[Dict[str, str]] = []

    def update(
        self,
        file_path: str,
        status: str,
        file_size: int = 0,
        error: Optional[str] = None,
    ):
        """Update progress."""
        self.processed_files += 1
        self.current_file = file_path
        self.bytes_processed += file_size

        if status == "success":
            self.successful += 1
        elif status == "error":
            self.failed += 1
            if error:
                self.errors.append({"file": file_path, "error": error})
        elif status == "skipped":
            self.skipped += 1

    def get_stats(self) -> Dict[str, Any]:
        """Get progress statistics."""
        elapsed = (datetime.now() - self.start_time).total_seconds()

        return {
            "total_files": self.total_files,
            "processed_files": self.processed_files,
            "successful": self.successful,
            "failed": self.failed,
            "skipped": self.skipped,
            "bytes_processed": self.bytes_processed,
            "current_file": self.current_file,
            "elapsed_seconds": elapsed,
            "files_per_second": self.processed_files / elapsed if elapsed > 0 else 0,
            "errors": self.errors[-10:],  # Last 10 errors
        }


class DocumentIndexer:
    """Index documents into the database."""

    # SQLite has a maximum blob size of 1 GB (1073741824 bytes)
    # We'll use a lower limit to be safe
    MAX_TEXT_SIZE = 500 * 1024 * 1024  # 500 MB

    def __init__(
        self,
        db: DatabaseManager,
        max_file_size: int = 400 * 1024 * 1024,
        trust_external_tools: bool = False,
        ignore_patterns: Optional[List[str]] = None,
        index_mode: str = "full",
    ):
        """
        Initialize document indexer.

        Args:
            db: Database manager
            max_file_size: Maximum file size to index (bytes)
            trust_external_tools: Whether to use external tools
            ignore_patterns: List of glob patterns to ignore
            index_mode: "full" for full-text indexing, "metadata_only" for just metadata
        """
        self.db = db
        self.max_file_size = max_file_size
        self.trust_external_tools = trust_external_tools
        self.ignore_patterns = ignore_patterns or []
        self.index_mode = index_mode
        self.extractor = TextExtractor(trust_external_tools) if index_mode == "full" else None
        self._stop_requested = False

    def index_directory(
        self,
        root_path: Path,
        reindex: bool = False,
        threads: int = 4,
        progress_callback: Optional[Callable[[IndexProgress], None]] = None,
        auto_mode_threshold: int = 400 * 1024 * 1024,
    ) -> IndexProgress:
        """
        Index all files in a directory.

        Args:
            root_path: Root directory to index
            reindex: Whether to reindex existing files
            threads: Number of worker threads
            progress_callback: Optional callback for progress updates
            auto_mode_threshold: Size threshold for auto mode (bytes)

        Returns:
            IndexProgress object with statistics
        """
        root_path = root_path.resolve()
        logger.info(f"Indexing directory: {root_path}")

        if not root_path.exists():
            raise ValueError(f"Path does not exist: {root_path}")

        if not root_path.is_dir():
            raise ValueError(f"Path is not a directory: {root_path}")

        # Collect files to index
        files = self._collect_files(root_path)

        # Auto-select index mode based on total size
        if self.index_mode == "auto":
            total_size = sum(f.stat().st_size for f in files)
            if total_size > auto_mode_threshold:
                actual_mode = "metadata_only"
                logger.info(
                    f"Auto mode: Total size {format_size(total_size)} > {format_size(auto_mode_threshold)}, "
                    f"using metadata_only mode"
                )
            else:
                actual_mode = "full"
                logger.info(
                    f"Auto mode: Total size {format_size(total_size)} <= {format_size(auto_mode_threshold)}, "
                    f"using full-text mode"
                )
            # Temporarily switch mode
            original_mode = self.index_mode
            self.index_mode = actual_mode
            if actual_mode == "full" and self.extractor is None:
                self.extractor = TextExtractor(self.trust_external_tools)
        else:
            original_mode = None

        progress = IndexProgress()
        progress.total_files = len(files)

        logger.info(f"Found {len(files)} files to index")

        # Warn if no files found
        if len(files) == 0:
            logger.warning(f"No files found to index in {root_path}. Check max_file_size setting or ignore patterns.")

        if reindex:
            logger.info("Removing existing entries for this root...")
            self.db.delete_by_root(str(root_path))
            existing_docs = {}
        else:
            # Batch-load existing documents for fast skip checking
            existing_docs = self.db.get_documents_metadata_by_root(str(root_path))

        # Process files
        if threads == 1:
            # Single-threaded
            for file_path in files:
                if self._stop_requested:
                    logger.info("Indexing stopped by user")
                    break
                self._index_file(file_path, root_path, progress, progress_callback, existing_docs)
        else:
            # Multi-threaded
            with concurrent.futures.ThreadPoolExecutor(max_workers=threads) as executor:
                futures = []
                for file_path in files:
                    if self._stop_requested:
                        logger.info("Indexing stopped by user")
                        break
                    future = executor.submit(
                        self._index_file,
                        file_path,
                        root_path,
                        progress,
                        progress_callback,
                        existing_docs,
                    )
                    futures.append(future)

                # Wait for all submitted futures
                concurrent.futures.wait(futures)

        logger.info(
            f"Indexing complete: {progress.successful} successful, {progress.failed} failed, {progress.skipped} skipped"
        )

        # Restore original mode if we switched in auto mode
        if original_mode is not None:
            self.index_mode = original_mode

        return progress

    def _collect_files(self, root_path: Path) -> List[Path]:
        """Collect all files to index."""
        files = []
        ignored_count = 0
        too_large_count = 0

        for path in root_path.rglob("*"):
            if not path.is_file():
                continue

            # Check ignore patterns
            if should_ignore(path, self.ignore_patterns):
                logger.debug(f"Ignoring {path} (matches ignore pattern)")
                ignored_count += 1
                continue

            # Check file size
            try:
                size = path.stat().st_size
                if size > self.max_file_size:
                    logger.debug(f"Ignoring {path} (too large: {format_size(size)})")
                    too_large_count += 1
                    continue
            except Exception as e:
                logger.warning(f"Failed to stat {path}: {e}")
                continue

            files.append(path)

        logger.info(f"Collected {len(files)} files to index (ignored: {ignored_count}, too large: {too_large_count})")
        return files

    def stop(self):
        """Request indexing to stop."""
        self._stop_requested = True
        logger.info("Stop requested for indexer")

    def _index_file(
        self,
        file_path: Path,
        root_path: Path,
        progress: IndexProgress,
        progress_callback: Optional[Callable[[IndexProgress], None]],
        existing_docs: Dict[str, Dict[str, Any]],
    ):
        """Index a single file."""
        # Check if stop was requested
        if self._stop_requested:
            return

        try:
            # Get file metadata
            stat = file_path.stat()
            file_size = stat.st_size
            mtime = stat.st_mtime

            # Detect file type
            file_type = detect_file_type(file_path)

            # Check if file needs reindexing (fast in-memory check)
            path_str = str(file_path)
            existing = existing_docs.get(path_str)
            if existing and existing.get("mtime") == mtime:
                logger.debug(f"Skipping {file_path} (not modified)")
                progress.update(path_str, "skipped", file_size)
                if progress_callback:
                    progress_callback(progress)
                return

            # Extract text (skip in metadata-only mode)
            extracted_text = ""
            extractor_name = "metadata-only"
            text_truncated = False

            if self.index_mode == "full":
                try:
                    extracted_text, extractor_name = self.extractor.extract(
                        file_path, file_type
                    )
                except Exception as e:
                    logger.error(f"Extraction failed for {file_path}: {e}")
                    extracted_text = ""
                    extractor_name = "error"

                # Check if extracted text is too large for SQLite
                if extracted_text and len(extracted_text) > self.MAX_TEXT_SIZE:
                    logger.warning(
                        f"Extracted text for {file_path} is too large ({len(extracted_text)} bytes), "
                        f"truncating to {self.MAX_TEXT_SIZE} bytes"
                    )
                    extracted_text = extracted_text[:self.MAX_TEXT_SIZE]
                    text_truncated = True

            # Calculate hash (skip in metadata-only mode for speed)
            sha256 = None
            if self.index_mode == "full":
                try:
                    sha256 = calculate_sha256(file_path)
                except Exception as e:
                    logger.warning(f"Failed to calculate hash for {file_path}: {e}")
                    sha256 = None

            # Insert into database
            if self.index_mode == "metadata_only":
                status = "success"
                error_message = "Metadata-only mode (use ripgrep for search)"
            else:
                status = "success" if extracted_text else "error"
                error_message = None
                if status == "error":
                    error_message = "No text extracted"
                elif text_truncated:
                    error_message = f"Text truncated (original size exceeded {self.MAX_TEXT_SIZE} bytes)"

            try:
                self.db.insert_document(
                    path=str(file_path),
                    root_path=str(root_path),
                    source_type=file_type,
                    extracted_text=extracted_text,
                    extractor=extractor_name,
                    sha256=sha256,
                    file_size=file_size,
                    mtime=mtime,
                    status=status,
                    error_message=error_message,
                )

                progress.update(str(file_path), status, file_size)
                logger.debug(f"Indexed {file_path} ({extractor_name})")

            except Exception as db_error:
                # Handle database errors (e.g., size limits)
                error_msg = str(db_error)
                logger.error(f"Database insert failed for {file_path}: {error_msg}")

                # If it's a size error and we haven't truncated yet, try with smaller text
                if "too big" in error_msg.lower() and not text_truncated:
                    logger.warning(f"Retrying with heavily truncated text for {file_path}")
                    extracted_text = extracted_text[:10 * 1024 * 1024]  # 10 MB only
                    try:
                        self.db.insert_document(
                            path=str(file_path),
                            root_path=str(root_path),
                            source_type=file_type,
                            extracted_text=extracted_text,
                            extractor=extractor_name,
                            sha256=sha256,
                            file_size=file_size,
                            mtime=mtime,
                            status="success",
                            error_message=f"Text heavily truncated due to database size limit",
                        )
                        progress.update(str(file_path), "success", file_size)
                        logger.info(f"Successfully indexed {file_path} with truncated text")
                    except Exception as retry_error:
                        logger.error(f"Retry failed for {file_path}: {retry_error}")
                        progress.update(str(file_path), "error", file_size, str(retry_error))
                else:
                    progress.update(str(file_path), "error", file_size, error_msg)

        except Exception as e:
            error_msg = f"{type(e).__name__}: {str(e)}"
            logger.error(f"Failed to index {file_path}: {error_msg}")
            logger.debug(traceback.format_exc())

            # Try to insert error record
            try:
                self.db.insert_document(
                    path=str(file_path),
                    root_path=str(root_path),
                    source_type=detect_file_type(file_path),
                    extracted_text="",
                    status="error",
                    error_message=error_msg,
                )
            except Exception as db_error:
                logger.error(f"Failed to record error for {file_path}: {db_error}")

            progress.update(str(file_path), "error", 0, error_msg)

        finally:
            if progress_callback:
                progress_callback(progress)

    def index_single_file(
        self, file_path: Path, root_path: Optional[Path] = None
    ) -> bool:
        """
        Index a single file.

        Args:
            file_path: Path to file
            root_path: Optional root path (defaults to file's parent)

        Returns:
            True if successful
        """
        if root_path is None:
            root_path = file_path.parent

        progress = IndexProgress()
        progress.total_files = 1

        # Use empty dict for existing_docs since we're indexing a single file
        self._index_file(file_path, root_path, progress, None, {})

        return progress.successful > 0
