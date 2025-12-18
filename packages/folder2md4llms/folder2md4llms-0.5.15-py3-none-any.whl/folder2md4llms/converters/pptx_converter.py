"""PowerPoint document converter."""

import logging
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from .base import BaseConverter, ConversionError

logger = logging.getLogger(__name__)


class PPTXConverter(BaseConverter):
    """Converts PPTX files to text."""

    def __init__(self, config: dict[str, Any] | None = None):
        super().__init__(config)
        self.max_slides = self.config.get("pptx_max_slides", 100)
        self.include_notes = self.config.get("pptx_include_notes", True)
        self.include_slide_numbers = self.config.get("pptx_include_slide_numbers", True)

    def can_convert(self, file_path: Path) -> bool:
        """Check if this converter can handle the given file."""
        return (
            PPTX_AVAILABLE
            and file_path.suffix.lower() in {".pptx", ".ppt"}
            and file_path.exists()
        )

    def convert(self, file_path: Path) -> str | None:
        """Convert PPTX to text."""
        if not PPTX_AVAILABLE:
            return "PowerPoint conversion not available. Install python-pptx: pip install python-pptx"

        try:
            # Load presentation
            prs = Presentation(str(file_path))

            # Get slide count
            slide_count = len(prs.slides)
            slides_to_process = min(slide_count, self.max_slides)

            # Format the output
            text_parts = []
            text_parts.append(f"PowerPoint Presentation: {file_path.name}")
            text_parts.append(f"Total slides: {slide_count}")

            if slides_to_process < slide_count:
                text_parts.append(f"Showing first {slides_to_process} slides")

            text_parts.append("=" * 50)
            text_parts.append("")

            # Process each slide
            for i, slide in enumerate(prs.slides[:slides_to_process]):
                if self.include_slide_numbers:
                    text_parts.append(f"--- Slide {i + 1} ---")
                    text_parts.append("")

                # Extract text from shapes
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    text_parts.append(slide_text)
                    text_parts.append("")

                # Extract notes if requested
                if self.include_notes and slide.has_notes_slide:
                    notes_text = self._extract_notes_text(slide)
                    if notes_text:
                        text_parts.append("*Speaker Notes:*")
                        text_parts.append(notes_text)
                        text_parts.append("")

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error converting PPTX {file_path}: {e}")
            raise ConversionError(f"Failed to convert PPTX: {str(e)}") from e

    def _extract_slide_text(self, slide) -> str:
        """Extract text from all shapes in a slide."""
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # Handle table content
                table_text = self._extract_table_text(shape)
                if table_text:
                    text_parts.append(table_text)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Handle grouped shapes
                group_text = self._extract_group_text(shape)
                if group_text:
                    text_parts.append(group_text)

        return "\n".join(text_parts)

    def _extract_table_text(self, table_shape) -> str:
        """Extract text from a table shape."""
        if not hasattr(table_shape, "table"):
            return ""

        table = table_shape.table
        table_parts = []

        for row in table.rows:
            row_parts = []
            for cell in row.cells:
                if hasattr(cell, "text") and cell.text.strip():
                    row_parts.append(cell.text.strip())
                else:
                    row_parts.append("")
            if any(row_parts):
                table_parts.append(" | ".join(row_parts))

        return "\n".join(table_parts)

    def _extract_group_text(self, group_shape) -> str:
        """Extract text from grouped shapes."""
        if not hasattr(group_shape, "shapes"):
            return ""

        text_parts = []
        for shape in group_shape.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        return "\n".join(text_parts)

    def _extract_notes_text(self, slide) -> str:
        """Extract text from slide notes."""
        try:
            notes_slide = slide.notes_slide
            if notes_slide and hasattr(notes_slide, "notes_text_frame"):
                text_frame = notes_slide.notes_text_frame
                if text_frame and hasattr(text_frame, "text"):
                    text = text_frame.text
                    return text.strip() if isinstance(text, str) else ""
        except Exception:
            return ""
        return ""

    def get_supported_extensions(self) -> set:
        """Get the file extensions this converter supports."""
        return {".pptx", ".ppt"}

    def get_document_info(self, file_path: Path) -> dict[str, Any]:
        """Get PPTX-specific information."""
        info = self.get_file_info(file_path)

        if not PPTX_AVAILABLE:
            info["error"] = "PowerPoint library not available"
            return info

        try:
            prs = Presentation(str(file_path))

            slide_count = len(prs.slides)
            slides_with_notes = sum(1 for slide in prs.slides if slide.has_notes_slide)

            info.update(
                {
                    "slides": slide_count,
                    "slides_with_notes": slides_with_notes,
                    "will_be_truncated": slide_count > self.max_slides,
                }
            )

            # Try to get presentation properties
            if hasattr(prs, "core_properties"):
                core_props = prs.core_properties
                info.update(
                    {
                        "title": getattr(core_props, "title", "") or "",
                        "author": getattr(core_props, "author", "") or "",
                        "subject": getattr(core_props, "subject", "") or "",
                        "created": str(getattr(core_props, "created", "")) or "",
                        "modified": str(getattr(core_props, "modified", "")) or "",
                    }
                )

        except Exception as e:
            info["error"] = str(e)

        return info
