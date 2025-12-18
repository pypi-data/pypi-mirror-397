import logging
import os
from typing import AsyncGenerator, Generator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docsloader import utils
from docsloader.base import BaseLoader, DocsData

logger = logging.getLogger(__name__)


class PptxLoader(BaseLoader):
    """
    pptx loader

    params:
        - path_or_url: str
        - suffix: str = None
        - encoding: str = None
        - load_type: str = "basic"
        - load_options: dict = None
        - metadata: dict = None
        - rm_tmpfile: bool = False
    """

    async def load_by_basic(self) -> AsyncGenerator[DocsData, None]:
        image_fmt = self.load_options.get("image_fmt")
        table_fmt = self.load_options.get("table_fmt")
        tmpfile_cvt = None
        if self.suffix == ".ppt":
            tmpfile_cvt = utils.office_cvt_openxml(filepath=self.tmpfile, file_suffix=self.suffix)
        for item in self.extract_by_python_pptx(
                filepath=tmpfile_cvt or self.tmpfile,
                image_fmt=image_fmt,
                table_fmt=table_fmt,
        ):
            self.metadata.update({
                "page": item.get("page"),
                "page_total": item.get("page_total"),
            })
            yield DocsData(
                type=item.get("type"),
                text=item.get("text"),
                data=item.get("data"),
                metadata=self.metadata,
            )
        if tmpfile_cvt:
            utils.rm_file(filepath=tmpfile_cvt)

    def extract_by_python_pptx(
            self,
            filepath: str,
            image_fmt: str,
            table_fmt: str,
    ) -> Generator[dict, None, None]:
        tmpdir = utils.mk_tmpdir()
        presentation = Presentation(filepath)
        page_total = len(presentation.slides)
        for slide_idx, slide in enumerate(presentation.slides):
            logger.debug(f"Processing slide {slide_idx + 1}")
            for shape_idx, shape in enumerate(slide.shapes):
                extracted_data = PptxLoader.extract_shape(
                    shape=shape,
                    tmpdir=tmpdir,
                    image_idx=f"{slide_idx}-{shape_idx}",
                    image_fmt=image_fmt,
                    table_fmt=table_fmt,
                )
                if extracted_data:
                    extracted_data.update(
                        page=slide_idx + 1,
                        page_total=page_total,
                    )
                    yield extracted_data
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub_shape_idx, sub_shape in enumerate(shape.shapes):
                        group_extracted_data = PptxLoader.extract_shape(
                            shape=sub_shape,
                            tmpdir=tmpdir,
                            image_idx=f"{slide_idx}-{shape_idx}-{sub_shape_idx}",
                            image_fmt=image_fmt,
                            table_fmt=table_fmt,
                        )
                        if group_extracted_data:
                            group_extracted_data.update(
                                page=slide_idx + 1,
                                page_total=page_total,
                            )
                            yield group_extracted_data
        utils.rm_empty_dir(tmpdir)

    @staticmethod
    def extract_shape(
            shape,
            tmpdir: str,
            image_idx: str,
            image_fmt: str = "path",
            table_fmt: str = "html",
    ) -> dict:
        """
        解析单个 shape 对象，提取其中的文本、表格和图片信息。
        """
        shape_text = ""
        shape_data = {}
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs).strip()
                if para_text:
                    shape_text += para_text + "\n"
            if shape_text:
                shape_data = {
                    "type": "text",
                    "text": shape_text,
                }
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_path = os.path.join(tmpdir, f"image_{image_idx}.{image.ext}")
            with open(image_path, "wb") as f:
                f.write(image.blob)
            shape_data = {
                "type": "image",
                "text": utils.format_image(image_path, fmt=image_fmt),  # noqa
                "data": image_path,
            }
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_title = shape.name if shape.name else "Table"
            table_data = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
            shape_text += f"\n## {table_title}\n"
            shape_text += utils.format_table(table_data, fmt=table_fmt)  # noqa
            shape_data = {
                "type": "table",
                "text": shape_text,
                "data": table_data
            }
        return shape_data
