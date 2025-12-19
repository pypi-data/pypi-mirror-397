import os
import sys
import tempfile
import traceback

import Orange
from Orange.data import StringVariable
from Orange.widgets import widget
from Orange.widgets.widget import Input, Output
from AnyQt.QtWidgets import QMessageBox, QApplication

from docx import Document
from docx.shared import Pt as pt_docx
from pptx import Presentation
from pptx.util import Inches, Pt
import pypandoc
from docx2pdf import convert

# Chargement UI
if "site-packages/Orange/widgets" in os.path.dirname(os.path.abspath(__file__)).replace("\\", "/"):
    from Orange.widgets.orangecontrib.AAIT.utils.import_uic import uic
    from Orange.widgets.orangecontrib.IO4IT.utils import utils_md
else:
    from orangecontrib.AAIT.utils.import_uic import uic
    from orangecontrib.IO4IT.utils import utils_md


class OWExportMarkdown(widget.OWWidget):
    name = "OWExportMarkdown"
    description = "Automatically export content to DOCX, PPTX, and PDF using the same base path."
    icon = "icons/export_md.png"
    if "site-packages/Orange/widgets" in os.path.dirname(os.path.abspath(__file__)).replace("\\", "/"):
        icon = "icons_dev/export_md.png"
    want_control_area = False
    priority = 9999
    category = "AAIT - TOOLBOX"

    class Inputs:
        data = Input("Data", Orange.data.Table)

    class Outputs:
        data = Output("Data", Orange.data.Table)

    def __init__(self):
        super().__init__()
        self.data = None
        ui_path = os.path.join(os.path.dirname(__file__), "designer", "owexportmarkdown.ui")
        uic.loadUi(ui_path, baseinstance=self)

    # -------- helpers headers/footers --------
    def ajouter_en_tete_pied_docx(self, file_path, header_text, footer_text):
        try:
            doc = Document(file_path)
            section = doc.sections[0]
            header = section.header
            p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
            p.text = header_text
            if p.runs:
                p.runs[0].font.size = pt_docx(10)
            footer = section.footer
            p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            p.text = footer_text
            if p.runs:
                p.runs[0].font.size = pt_docx(10)
            doc.save(file_path)
        except Exception:
            # en-tête/pied non bloquants
            pass

    def ajouter_entete_pied_pptx(self, file_path, entete_text, pied_text):
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                entete = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(8), Inches(0.5))
                tf_entete = entete.text_frame
                tf_entete.text = entete_text
                tf_entete.paragraphs[0].font.size = Pt(12)
                tf_entete.paragraphs[0].font.bold = True

                pied = slide.shapes.add_textbox(Inches(0.3), Inches(6.3), Inches(8), Inches(0.5))
                tf_pied = pied.text_frame
                tf_pied.text = pied_text
                tf_pied.paragraphs[0].font.size = Pt(10)
            prs.save(file_path)
        except Exception:
            pass

    # -------------- input --------------
    @Inputs.data
    def set_data(self, in_data):
        self.error("")
        if in_data is None:
            self.data = None
            self.Outputs.data.send(None)
            return

        # On exige au moins 'path'
        if "path" not in in_data.domain:
            self.error("La table d'entrée doit contenir au moins la colonne 'path'.")
            self.Outputs.data.send(None)
            return

        # Optionnellement 'content'
        self.data = in_data
        try:
            table_out = self.export_all_rows()
            self.Outputs.data.send(table_out)
        except Exception as e:
            tb = traceback.format_exc()
            QMessageBox.critical(self, "Erreur d'export", f"{e}\n\n{tb}")
            self.Outputs.data.send(None)

    # -------------- core --------------
    def export_all_rows(self):
        base_paths = self.data.get_column("path")
        has_content = "content" in self.data.domain
        file_contents = self.data.get_column("content") if has_content else [None] * len(base_paths)

        pdf_paths, docx_paths, pptx_paths = [], [], []

        for i, (md_text, base_path) in enumerate(zip(file_contents, base_paths)):
            base_path = str(base_path or "").strip()

            # Lecture du contenu si 'content' absent et path en .md
            if not has_content:
                if base_path.lower().endswith(".md"):
                    try:
                        with open(base_path, "r", encoding="utf-8") as f:
                            md_text = f.read()
                        # on remplace base_path par le même (on garde la base pour sorties)
                    except Exception as e:
                        self.error(f"Impossible de lire le fichier : {base_path} ({e})")
                        pdf_paths.append("")
                        docx_paths.append("")
                        pptx_paths.append("")
                        continue
                else:
                    # pas de content et path non .md -> rien à faire pour cette ligne
                    pdf_paths.append("")
                    docx_paths.append("")
                    pptx_paths.append("")
                    continue

            md_text = (str(md_text or "")).strip()

            if not md_text or not base_path:
                pdf_paths.append("")
                docx_paths.append("")
                pptx_paths.append("")
                continue

            # Normaliser la base: enlever extension si présente
            base_no_ext, _ = os.path.splitext(base_path)
            # Créer dossier si nécessaire
            out_dir = os.path.dirname(base_no_ext)
            if out_dir and not os.path.isdir(out_dir):
                os.makedirs(out_dir, exist_ok=True)

            docx_out = base_no_ext + ".docx"
            pptx_out = base_no_ext + ".pptx"
            pdf_out = base_no_ext + ".pdf"

            # MD temporaire
            with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as tmp:
                tmp.write(md_text)
                tmp_md = tmp.name

            try:
                if utils_md.is_word_installed():
                    pypandoc.convert_file(tmp_md, to="docx", format="gfm-yaml_metadata_block", outputfile=docx_out)
                    self.ajouter_en_tete_pied_docx(
                        docx_out,
                        "Rapport - Orange AI",
                        "Page générée automatiquement - Ne pas diffuser"
                    )

                    # PPTX
                    pypandoc.convert_file(tmp_md, to="pptx", format="gfm-yaml_metadata_block", outputfile=pptx_out)
                    self.ajouter_entete_pied_pptx(
                        pptx_out,
                        "Orange AI – Présentation",
                        "Page générée automatiquement"
                    )
                else:
                    raise Exception("Word non détecté")

            except Exception:
                try:
                    pypandoc.convert_file(tmp_md, to="pdf", outputfile=pdf_out)
                except Exception:
                    self.error(f"Échec conversion PDF pour la ligne {i + 1}.")
                    pdf_out = ""
            finally:
                try:
                    os.remove(tmp_md)
                except Exception:
                    pass

            pdf_paths.append(pdf_out if os.path.isfile(pdf_out) else "")
            docx_paths.append(docx_out if os.path.isfile(docx_out) else "")
            pptx_paths.append(pptx_out if os.path.isfile(pptx_out) else "")

        # Ajouter colonnes sortie
        table = self.data
        table = table.add_column(StringVariable("output_pdf_path"), pdf_paths)
        table = table.add_column(StringVariable("output_docx_path"), docx_paths)
        table = table.add_column(StringVariable("output_pptx_path"), pptx_paths)

        return table


if __name__ == "__main__":
    app = QApplication(sys.argv)
    w = OWExportMarkdown()
    w.show()
    sys.exit(app.exec() if hasattr(app, "exec") else app.exec_())
