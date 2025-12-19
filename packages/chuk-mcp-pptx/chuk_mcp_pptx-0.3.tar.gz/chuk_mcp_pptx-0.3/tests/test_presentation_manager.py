"""
Comprehensive tests for PresentationManager.

Tests cover:
- Initialization and configuration
- CRUD operations (create, read, update, delete)
- Artifact store integration
- Base64 import/export
- Metadata management
- Slide metadata updates
- Error handling

All tests use:
- Pydantic models (no dict goop)
- Async/await patterns
- Design tokens where applicable
- Proper type hints
"""

import pytest
from unittest.mock import MagicMock, AsyncMock, patch
from datetime import datetime

from chuk_mcp_pptx.core.presentation_manager import PresentationManager
from chuk_mcp_pptx.models import (
    PresentationMetadata,
    PresentationInfo,
    ListPresentationsResponse,
)


class TestPresentationManagerInitialization:
    """Tests for PresentationManager initialization."""

    def test_init_default_base_path(self) -> None:
        """Test initialization with default base path."""
        manager = PresentationManager()
        assert manager.base_path == "presentations"

    def test_init_custom_base_path(self) -> None:
        """Test initialization with custom base path."""
        manager = PresentationManager(base_path="custom/path")
        assert manager.base_path == "custom/path"

    def test_init_empty_collections(self) -> None:
        """Test that collections are initialized empty."""
        manager = PresentationManager()
        assert len(manager._presentations) == 0
        assert len(manager._metadata) == 0
        assert len(manager._namespace_ids) == 0
        assert manager._current_presentation is None

    def test_mime_type_constant(self) -> None:
        """Test PPTX MIME type constant is correct."""
        assert PresentationManager.PPTX_MIME_TYPE == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )


class TestSanitizeName:
    """Tests for name sanitization."""

    def test_sanitize_alphanumeric(self) -> None:
        """Test sanitization preserves alphanumeric characters."""
        manager = PresentationManager()
        assert manager._sanitize_name("test123") == "test123"

    def test_sanitize_with_dashes(self) -> None:
        """Test sanitization preserves dashes."""
        manager = PresentationManager()
        assert manager._sanitize_name("my-presentation") == "my-presentation"

    def test_sanitize_with_underscores(self) -> None:
        """Test sanitization preserves underscores."""
        manager = PresentationManager()
        assert manager._sanitize_name("my_presentation") == "my_presentation"

    def test_sanitize_removes_special_chars(self) -> None:
        """Test sanitization removes special characters."""
        manager = PresentationManager()
        assert manager._sanitize_name("test@#$%file") == "testfile"

    def test_sanitize_removes_spaces(self) -> None:
        """Test sanitization removes spaces."""
        manager = PresentationManager()
        assert manager._sanitize_name("my presentation") == "mypresentation"

    def test_sanitize_empty_result_fallback(self) -> None:
        """Test sanitization falls back to 'presentation' for empty result."""
        manager = PresentationManager()
        assert manager._sanitize_name("@#$%") == "presentation"

    def test_sanitize_empty_string(self) -> None:
        """Test sanitization of empty string."""
        manager = PresentationManager()
        assert manager._sanitize_name("") == "presentation"


class TestGetStore:
    """Tests for artifact store retrieval."""

    def test_get_store_when_available(self) -> None:
        """Test getting store when artifact store is available."""
        manager = PresentationManager()
        mock_store = MagicMock()

        # Mock at the chuk_mcp_server module level since import happens inside method
        with patch("chuk_mcp_server.has_artifact_store", return_value=True):
            with patch("chuk_mcp_server.get_artifact_store", return_value=mock_store):
                store = manager._get_store()
                assert store == mock_store

    def test_get_store_when_not_available(self) -> None:
        """Test getting store when artifact store is not available."""
        manager = PresentationManager()

        # Mock at the chuk_mcp_server module level since import happens inside method
        with patch("chuk_mcp_server.has_artifact_store", return_value=False):
            store = manager._get_store()
            assert store is None


class TestNamespaceAndArtifactUri:
    """Tests for namespace ID and artifact URI methods."""

    def test_get_namespace_id_exists(self) -> None:
        """Test getting namespace ID when it exists."""
        manager = PresentationManager()
        manager._namespace_ids["test"] = "ns-123"
        assert manager.get_namespace_id("test") == "ns-123"

    def test_get_namespace_id_not_exists(self) -> None:
        """Test getting namespace ID when it doesn't exist."""
        manager = PresentationManager()
        assert manager.get_namespace_id("nonexistent") is None

    def test_get_artifact_uri_exists(self) -> None:
        """Test getting artifact URI when namespace exists."""
        manager = PresentationManager(base_path="presentations")
        manager._namespace_ids["test"] = "ns-123"
        uri = manager.get_artifact_uri("test")
        assert uri == "artifact://chuk-mcp-pptx/presentations/test"

    def test_get_artifact_uri_not_exists(self) -> None:
        """Test getting artifact URI when namespace doesn't exist."""
        manager = PresentationManager()
        assert manager.get_artifact_uri("nonexistent") is None


class TestCreatePresentation:
    """Tests for presentation creation."""

    @pytest.mark.asyncio
    async def test_create_basic(self) -> None:
        """Test basic presentation creation."""
        manager = PresentationManager()
        metadata = await manager.create(name="test_presentation")

        assert isinstance(metadata, PresentationMetadata)
        assert metadata.name == "test_presentation"
        assert metadata.slide_count == 0
        assert metadata.theme is None
        assert "test_presentation" in manager._presentations
        assert manager._current_presentation == "test_presentation"

    @pytest.mark.asyncio
    async def test_create_with_theme(self) -> None:
        """Test presentation creation with theme."""
        manager = PresentationManager()
        metadata = await manager.create(name="themed_presentation", theme="dark")

        assert metadata.theme == "dark"

    @pytest.mark.asyncio
    async def test_create_sets_current(self) -> None:
        """Test that create sets the current presentation."""
        manager = PresentationManager()
        await manager.create(name="first")
        assert manager._current_presentation == "first"

        await manager.create(name="second")
        assert manager._current_presentation == "second"

    @pytest.mark.asyncio
    async def test_create_metadata_stored(self) -> None:
        """Test that metadata is stored properly."""
        manager = PresentationManager()
        await manager.create(name="test")

        assert "test" in manager._metadata
        metadata = manager._metadata["test"]
        assert isinstance(metadata, PresentationMetadata)
        assert isinstance(metadata.created_at, datetime)
        assert isinstance(metadata.modified_at, datetime)


class TestGetPresentation:
    """Tests for getting presentations."""

    @pytest.mark.asyncio
    async def test_get_existing_presentation(self) -> None:
        """Test getting an existing presentation."""
        manager = PresentationManager()
        await manager.create(name="test")

        result = await manager.get(name="test")
        assert result is not None
        prs, metadata = result
        assert metadata.name == "test"

    @pytest.mark.asyncio
    async def test_get_current_presentation(self) -> None:
        """Test getting the current presentation without specifying name."""
        manager = PresentationManager()
        await manager.create(name="test")

        result = await manager.get()
        assert result is not None
        _, metadata = result
        assert metadata.name == "test"

    @pytest.mark.asyncio
    async def test_get_nonexistent_presentation(self) -> None:
        """Test getting a nonexistent presentation."""
        manager = PresentationManager()
        result = await manager.get(name="nonexistent")
        assert result is None

    @pytest.mark.asyncio
    async def test_get_no_current_presentation(self) -> None:
        """Test getting when no current presentation is set."""
        manager = PresentationManager()
        result = await manager.get()
        assert result is None

    @pytest.mark.asyncio
    async def test_get_creates_metadata_if_missing(self) -> None:
        """Test that get creates metadata if it's missing."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Remove metadata manually
        del manager._metadata["test"]

        result = await manager.get(name="test")
        assert result is not None
        _, metadata = result
        assert metadata.name == "test"
        assert "test" in manager._metadata


class TestGetPresentationAsync:
    """Tests for async presentation retrieval."""

    @pytest.mark.asyncio
    async def test_get_presentation_exists(self) -> None:
        """Test getting presentation object when it exists."""
        manager = PresentationManager()
        await manager.create(name="test")

        prs = await manager.get_presentation(name="test")
        assert prs is not None

    @pytest.mark.asyncio
    async def test_get_presentation_not_exists(self) -> None:
        """Test getting presentation object when it doesn't exist."""
        manager = PresentationManager()
        prs = await manager.get_presentation(name="nonexistent")
        assert prs is None

    @pytest.mark.asyncio
    async def test_get_presentation_current(self) -> None:
        """Test getting current presentation without name."""
        manager = PresentationManager()
        await manager.create(name="test")

        prs = await manager.get_presentation()
        assert prs is not None

    @pytest.mark.asyncio
    async def test_get_presentation_no_current(self) -> None:
        """Test getting presentation when no current is set."""
        manager = PresentationManager()
        prs = await manager.get_presentation()
        assert prs is None


class TestGetMetadata:
    """Tests for metadata retrieval."""

    @pytest.mark.asyncio
    async def test_get_metadata_exists(self) -> None:
        """Test getting metadata when it exists."""
        manager = PresentationManager()
        await manager.create(name="test")

        metadata = await manager.get_metadata(name="test")
        assert metadata is not None
        assert isinstance(metadata, PresentationMetadata)
        assert metadata.name == "test"

    @pytest.mark.asyncio
    async def test_get_metadata_not_exists(self) -> None:
        """Test getting metadata when it doesn't exist."""
        manager = PresentationManager()
        metadata = await manager.get_metadata(name="nonexistent")
        assert metadata is None

    @pytest.mark.asyncio
    async def test_get_metadata_current(self) -> None:
        """Test getting current presentation metadata without name."""
        manager = PresentationManager()
        await manager.create(name="test")

        metadata = await manager.get_metadata()
        assert metadata is not None
        assert metadata.name == "test"

    @pytest.mark.asyncio
    async def test_get_metadata_no_current(self) -> None:
        """Test getting metadata when no current is set."""
        manager = PresentationManager()
        metadata = await manager.get_metadata()
        assert metadata is None


class TestSavePresentation:
    """Tests for saving presentations."""

    @pytest.mark.asyncio
    async def test_save_existing(self) -> None:
        """Test saving an existing presentation."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Save should work (returns False without artifact store)
        result = await manager.save(name="test")
        # Result depends on artifact store availability
        assert isinstance(result, bool)

    @pytest.mark.asyncio
    async def test_save_nonexistent(self) -> None:
        """Test saving a nonexistent presentation."""
        manager = PresentationManager()
        result = await manager.save(name="nonexistent")
        assert result is False

    @pytest.mark.asyncio
    async def test_save_current(self) -> None:
        """Test saving current presentation without specifying name."""
        manager = PresentationManager()
        await manager.create(name="test")

        result = await manager.save()
        assert isinstance(result, bool)


class TestUpdatePresentation:
    """Tests for updating presentations."""

    @pytest.mark.asyncio
    async def test_update_existing(self) -> None:
        """Test updating an existing presentation."""
        manager = PresentationManager()
        await manager.create(name="test")

        original_modified = manager._metadata["test"].modified_at

        # Small delay to ensure time difference
        import asyncio

        await asyncio.sleep(0.01)

        result = await manager.update(name="test")
        assert isinstance(result, bool)

        # Modified time should be updated
        new_modified = manager._metadata["test"].modified_at
        assert new_modified >= original_modified

    @pytest.mark.asyncio
    async def test_update_nonexistent(self) -> None:
        """Test updating a nonexistent presentation."""
        manager = PresentationManager()
        result = await manager.update(name="nonexistent")
        assert result is False

    @pytest.mark.asyncio
    async def test_update_updates_slide_count(self) -> None:
        """Test that update refreshes slide count in metadata."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Get presentation and add a slide
        prs = manager._presentations["test"]
        prs.slides.add_slide(prs.slide_layouts[0])

        await manager.update(name="test")

        metadata = manager._metadata["test"]
        assert metadata.slide_count == 1


class TestDeletePresentation:
    """Tests for deleting presentations."""

    @pytest.mark.asyncio
    async def test_delete_existing(self) -> None:
        """Test deleting an existing presentation."""
        manager = PresentationManager()
        await manager.create(name="test")

        result = await manager.delete(name="test")
        assert result is True
        assert "test" not in manager._presentations
        assert "test" not in manager._metadata

    @pytest.mark.asyncio
    async def test_delete_nonexistent(self) -> None:
        """Test deleting a nonexistent presentation."""
        manager = PresentationManager()
        result = await manager.delete(name="nonexistent")
        assert result is False

    @pytest.mark.asyncio
    async def test_delete_updates_current(self) -> None:
        """Test that deleting current presentation updates current."""
        manager = PresentationManager()
        await manager.create(name="first")
        await manager.create(name="second")

        # Current is now "second"
        assert manager._current_presentation == "second"

        await manager.delete(name="second")

        # Current should be updated to "first"
        assert manager._current_presentation == "first"

    @pytest.mark.asyncio
    async def test_delete_last_presentation(self) -> None:
        """Test deleting the last presentation."""
        manager = PresentationManager()
        await manager.create(name="only")

        await manager.delete(name="only")

        assert manager._current_presentation is None
        assert len(manager._presentations) == 0


class TestListPresentations:
    """Tests for listing presentations."""

    @pytest.mark.asyncio
    async def test_list_empty(self) -> None:
        """Test listing when no presentations exist."""
        manager = PresentationManager()
        response = await manager.list_presentations()

        assert isinstance(response, ListPresentationsResponse)
        assert response.total == 0
        assert len(response.presentations) == 0
        assert response.current is None

    @pytest.mark.asyncio
    async def test_list_single(self) -> None:
        """Test listing with a single presentation."""
        manager = PresentationManager()
        await manager.create(name="test")

        response = await manager.list_presentations()

        assert response.total == 1
        assert len(response.presentations) == 1
        assert response.current == "test"

        pres_info = response.presentations[0]
        assert isinstance(pres_info, PresentationInfo)
        assert pres_info.name == "test"
        assert pres_info.is_current is True

    @pytest.mark.asyncio
    async def test_list_multiple(self) -> None:
        """Test listing with multiple presentations."""
        manager = PresentationManager()
        await manager.create(name="first")
        await manager.create(name="second")

        response = await manager.list_presentations()

        assert response.total == 2
        assert len(response.presentations) == 2

        names = {p.name for p in response.presentations}
        assert names == {"first", "second"}

        # Only current should be marked as current
        current_count = sum(1 for p in response.presentations if p.is_current)
        assert current_count == 1


class TestSetCurrent:
    """Tests for setting current presentation."""

    @pytest.mark.asyncio
    async def test_set_current_existing(self) -> None:
        """Test setting current to an existing presentation."""
        manager = PresentationManager()
        await manager.create(name="first")
        await manager.create(name="second")

        result = await manager.set_current(name="first")
        assert result is True
        assert manager._current_presentation == "first"

    @pytest.mark.asyncio
    async def test_set_current_nonexistent(self) -> None:
        """Test setting current to a nonexistent presentation."""
        manager = PresentationManager()
        result = await manager.set_current(name="nonexistent")
        assert result is False

    def test_get_current_name(self) -> None:
        """Test getting current presentation name."""
        manager = PresentationManager()
        assert manager.get_current_name() is None

    @pytest.mark.asyncio
    async def test_get_current_name_after_create(self) -> None:
        """Test getting current name after creating a presentation."""
        manager = PresentationManager()
        await manager.create(name="test")
        assert manager.get_current_name() == "test"


class TestUpdateSlideMetadata:
    """Tests for updating slide metadata."""

    @pytest.mark.asyncio
    async def test_update_slide_metadata_basic(self) -> None:
        """Test basic slide metadata update."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Add a slide
        prs = manager._presentations["test"]
        prs.slides.add_slide(prs.slide_layouts[0])

        await manager.update_slide_metadata(slide_index=0)

        metadata = manager._metadata["test"]
        assert len(metadata.slides) >= 1
        assert metadata.slides[0].index == 0

    @pytest.mark.asyncio
    async def test_update_slide_metadata_no_current(self) -> None:
        """Test slide metadata update with no current presentation."""
        manager = PresentationManager()
        # Should not raise
        await manager.update_slide_metadata(slide_index=0)

    @pytest.mark.asyncio
    async def test_update_slide_metadata_expands_list(self) -> None:
        """Test that slide metadata list expands as needed."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Add multiple slides
        prs = manager._presentations["test"]
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[0])

        await manager.update_slide_metadata(slide_index=2)

        metadata = manager._metadata["test"]
        assert len(metadata.slides) >= 3


class TestExportBase64:
    """Tests for base64 export."""

    @pytest.mark.asyncio
    async def test_export_basic(self) -> None:
        """Test basic base64 export."""
        manager = PresentationManager()
        await manager.create(name="test")

        result = await manager.export_base64(name="test")
        assert result is not None
        assert isinstance(result, str)
        # Should be valid base64
        import base64

        decoded = base64.b64decode(result)
        assert len(decoded) > 0

    @pytest.mark.asyncio
    async def test_export_nonexistent(self) -> None:
        """Test export of nonexistent presentation."""
        manager = PresentationManager()
        result = await manager.export_base64(name="nonexistent")
        assert result is None

    @pytest.mark.asyncio
    async def test_export_current(self) -> None:
        """Test export of current presentation without name."""
        manager = PresentationManager()
        await manager.create(name="test")

        result = await manager.export_base64()
        assert result is not None


class TestImportBase64:
    """Tests for base64 import."""

    @pytest.mark.asyncio
    async def test_import_from_export(self) -> None:
        """Test importing a presentation that was exported."""
        manager = PresentationManager()
        await manager.create(name="original")

        # Add a slide
        prs = manager._presentations["original"]
        prs.slides.add_slide(prs.slide_layouts[0])

        # Export
        exported = await manager.export_base64(name="original")
        assert exported is not None

        # Import with new name
        result = await manager.import_base64(data=exported, name="imported")
        assert result is True
        assert "imported" in manager._presentations
        assert manager._current_presentation == "imported"

    @pytest.mark.asyncio
    async def test_import_creates_metadata(self) -> None:
        """Test that import creates proper metadata."""
        manager = PresentationManager()
        await manager.create(name="original")

        exported = await manager.export_base64(name="original")
        assert exported is not None

        await manager.import_base64(data=exported, name="imported")

        metadata = manager._metadata.get("imported")
        assert metadata is not None
        assert isinstance(metadata, PresentationMetadata)
        assert metadata.name == "imported"

    @pytest.mark.asyncio
    async def test_import_invalid_base64(self) -> None:
        """Test importing invalid base64 data."""
        manager = PresentationManager()
        result = await manager.import_base64(data="not-valid-base64!", name="test")
        assert result is False


class TestClearAll:
    """Tests for clearing all presentations."""

    @pytest.mark.asyncio
    async def test_clear_all(self) -> None:
        """Test clearing all presentations."""
        manager = PresentationManager()
        await manager.create(name="first")
        await manager.create(name="second")

        manager.clear_all()

        assert len(manager._presentations) == 0
        assert len(manager._metadata) == 0
        assert len(manager._namespace_ids) == 0
        assert manager._current_presentation is None


class TestArtifactStoreIntegration:
    """Tests for artifact store integration."""

    @pytest.mark.asyncio
    async def test_save_to_store_no_store(self) -> None:
        """Test saving when no artifact store is available."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Without artifact store configured, save returns False
        with patch.object(manager, "_get_store", return_value=None):
            result = await manager._save_to_store("test", manager._presentations["test"])
            assert result is False

    @pytest.mark.asyncio
    async def test_load_from_store_no_store(self) -> None:
        """Test loading when no artifact store is available."""
        manager = PresentationManager()

        with patch.object(manager, "_get_store", return_value=None):
            result = await manager._load_from_store("test")
            assert result is None

    @pytest.mark.asyncio
    async def test_delete_from_store_no_store(self) -> None:
        """Test deleting when no artifact store is available."""
        manager = PresentationManager()

        with patch.object(manager, "_get_store", return_value=None):
            result = await manager._delete_from_store("test")
            assert result is False

    @pytest.mark.asyncio
    async def test_delete_from_store_no_namespace(self) -> None:
        """Test deleting when namespace doesn't exist."""
        manager = PresentationManager()
        mock_store = MagicMock()

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._delete_from_store("nonexistent")
            assert result is False

    @pytest.mark.asyncio
    async def test_save_to_store_new_namespace(self) -> None:
        """Test saving to store creates new namespace."""
        from pptx import Presentation as PptxPresentation

        manager = PresentationManager()
        prs = PptxPresentation()

        # Create mock store with async methods
        mock_store = MagicMock()
        mock_namespace_info = MagicMock()
        mock_namespace_info.namespace_id = "ns-test-123"
        mock_store.create_namespace = AsyncMock(return_value=mock_namespace_info)
        mock_store.write_namespace = AsyncMock(return_value=None)

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._save_to_store("new_pres", prs)

        assert result is True
        assert manager._namespace_ids["new_pres"] == "ns-test-123"
        mock_store.create_namespace.assert_called_once()
        mock_store.write_namespace.assert_called_once()

    @pytest.mark.asyncio
    async def test_save_to_store_existing_namespace(self) -> None:
        """Test saving to store with existing namespace updates it."""
        from pptx import Presentation as PptxPresentation

        manager = PresentationManager()
        prs = PptxPresentation()
        manager._namespace_ids["existing_pres"] = "ns-existing-456"

        # Create mock store with async methods
        mock_store = MagicMock()
        mock_store.write_namespace = AsyncMock(return_value=None)

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._save_to_store("existing_pres", prs)

        assert result is True
        # Should not call create_namespace for existing
        assert not hasattr(mock_store, "create_namespace") or not mock_store.create_namespace.called
        mock_store.write_namespace.assert_called_once()

    @pytest.mark.asyncio
    async def test_save_to_store_exception_handling(self) -> None:
        """Test that save handles exceptions gracefully."""
        from pptx import Presentation as PptxPresentation

        manager = PresentationManager()
        prs = PptxPresentation()

        mock_store = MagicMock()
        mock_store.create_namespace = AsyncMock(side_effect=Exception("Storage error"))

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._save_to_store("error_pres", prs)

        assert result is False

    @pytest.mark.asyncio
    async def test_load_from_store_success(self) -> None:
        """Test loading presentation from store successfully."""
        from pptx import Presentation as PptxPresentation
        import io

        manager = PresentationManager()

        # Create actual PPTX data
        prs = PptxPresentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_data = buffer.read()

        # Create mock namespace info
        mock_ns_info = MagicMock()
        mock_ns_info.name = "presentations/stored_pres"
        mock_ns_info.namespace_id = "ns-stored-789"

        mock_store = MagicMock()
        mock_store.list_namespaces = AsyncMock(return_value=[mock_ns_info])
        mock_store.read_namespace = AsyncMock(return_value=pptx_data)

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._load_from_store("stored_pres")

        assert result is not None
        mock_store.read_namespace.assert_called_once_with("ns-stored-789")

    @pytest.mark.asyncio
    async def test_load_from_store_no_namespace_id(self) -> None:
        """Test loading when namespace ID doesn't exist."""

        manager = PresentationManager()

        mock_store = MagicMock()
        mock_store.read_namespace = AsyncMock()

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._load_from_store("unknown_pres")

        assert result is None
        mock_store.read_namespace.assert_not_called()

    @pytest.mark.asyncio
    async def test_load_from_store_no_data(self) -> None:
        """Test loading when store returns None."""

        manager = PresentationManager()
        manager._namespace_ids["empty_pres"] = "ns-empty-000"

        mock_store = MagicMock()
        mock_store.read_namespace = AsyncMock(return_value=None)

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._load_from_store("empty_pres")

        assert result is None

    @pytest.mark.asyncio
    async def test_load_from_store_exception_handling(self) -> None:
        """Test that load handles exceptions gracefully."""

        manager = PresentationManager()
        manager._namespace_ids["error_pres"] = "ns-error-111"

        mock_store = MagicMock()
        mock_store.read_namespace = AsyncMock(side_effect=Exception("Read error"))

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._load_from_store("error_pres")

        assert result is None

    @pytest.mark.asyncio
    async def test_delete_from_store_success(self) -> None:
        """Test deleting presentation from store successfully."""

        manager = PresentationManager()
        manager._namespace_ids["delete_pres"] = "ns-delete-222"

        mock_store = MagicMock()
        mock_store.destroy_namespace = AsyncMock(return_value=None)

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._delete_from_store("delete_pres")

        assert result is True
        assert "delete_pres" not in manager._namespace_ids
        mock_store.destroy_namespace.assert_called_once_with("ns-delete-222")

    @pytest.mark.asyncio
    async def test_delete_from_store_exception_handling(self) -> None:
        """Test that delete handles exceptions gracefully."""

        manager = PresentationManager()
        manager._namespace_ids["error_del"] = "ns-error-del-333"

        mock_store = MagicMock()
        mock_store.destroy_namespace = AsyncMock(side_effect=Exception("Delete error"))

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager._delete_from_store("error_del")

        assert result is False

    @pytest.mark.asyncio
    async def test_get_loads_from_store_on_miss(self) -> None:
        """Test that get() loads from store when not in memory."""
        from pptx import Presentation as PptxPresentation
        import io

        manager = PresentationManager()

        # Create actual PPTX data
        prs = PptxPresentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_data = buffer.read()

        # Create mock namespace info
        mock_ns_info = MagicMock()
        mock_ns_info.name = "presentations/external_pres"
        mock_ns_info.namespace_id = "ns-external-444"

        mock_store = MagicMock()
        mock_store.list_namespaces = AsyncMock(return_value=[mock_ns_info])
        mock_store.read_namespace = AsyncMock(return_value=pptx_data)

        with patch.object(manager, "_get_store", return_value=mock_store):
            result = await manager.get("external_pres")

        assert result is not None
        prs_obj, metadata = result
        assert metadata.name == "external_pres"
        assert "external_pres" in manager._presentations


class TestEdgeCases:
    """Tests for edge cases and error handling."""

    @pytest.mark.asyncio
    async def test_create_overwrite_existing(self) -> None:
        """Test creating a presentation with a name that already exists."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Creating again should overwrite
        await manager.create(name="test")

        # Should still have one presentation
        assert len(manager._presentations) == 1

    @pytest.mark.asyncio
    async def test_metadata_timestamps(self) -> None:
        """Test that timestamps are properly set."""
        manager = PresentationManager()

        before = datetime.now()
        await manager.create(name="test")
        after = datetime.now()

        metadata = manager._metadata["test"]
        assert before <= metadata.created_at <= after
        assert before <= metadata.modified_at <= after

    @pytest.mark.asyncio
    async def test_presentation_not_in_metadata(self) -> None:
        """Test handling when presentation exists but metadata doesn't - should create metadata."""
        from unittest.mock import patch

        manager = PresentationManager()
        # This shouldn't happen normally, but test the handling
        from pptx import Presentation

        manager._presentations["orphan"] = Presentation()
        manager._metadata.pop("orphan", None)  # Ensure no metadata
        manager._current_presentation = "orphan"
        manager._cache_timestamps["orphan"] = float("inf")  # Mark cache as valid

        # Mock the store to return None (no artifact store available)
        with patch.object(manager, "_get_store", return_value=None):
            # get_metadata should now create metadata if presentation exists in cache
            metadata = await manager.get_metadata(name="orphan")
            assert metadata is not None  # Metadata is created for existing presentations
            assert metadata.name == "orphan"

    @pytest.mark.asyncio
    async def test_slide_metadata_index_bounds(self) -> None:
        """Test slide metadata update with out-of-bounds index."""
        manager = PresentationManager()
        await manager.create(name="test")

        # Update with high index - should expand list
        await manager.update_slide_metadata(slide_index=10)

        metadata = manager._metadata["test"]
        # List should be expanded but slide won't exist
        assert len(metadata.slides) >= 11
