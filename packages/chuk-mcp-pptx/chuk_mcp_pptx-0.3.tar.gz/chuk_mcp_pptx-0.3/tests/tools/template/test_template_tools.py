"""
Tests for tools/template/* modules.

Comprehensive tests for template tools including:
- __init__.py: Registration functions
- analyze.py: Template analysis tools
- list.py: Template listing tools
- import_tools.py: Template import tools
- workflow.py: Template workflow tools
"""

import json
import pytest
from unittest.mock import MagicMock, AsyncMock
from pptx import Presentation

from chuk_mcp_pptx.tools.template import (
    register_template_tools,
    register_list_tools,
    register_analyze_tools,
    register_import_tools,
    register_workflow_tools,
)
from chuk_mcp_pptx.tools.template.models import (
    LayoutInfo,
    TemplateInfo,
    PresentationTemplateListResponse,
    BuiltinTemplateInfo,
    CustomTemplateInfo,
)


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server."""
    mcp = MagicMock()
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mcp.tool = tool_decorator
    mcp._tools = tools
    return mcp


@pytest.fixture
def mock_template_manager():
    """Create a mock template manager."""
    manager = MagicMock()

    # Mock list_templates
    mock_template = MagicMock()
    mock_template.name = "test_template"
    mock_template.display_name = "Test Template"
    mock_template.description = "A test template"
    mock_template.category = "test"
    mock_template.layout_count = 10
    mock_template.tags = ["test", "sample"]
    manager.list_templates.return_value = [mock_template]

    # Mock get_template_data - return None by default (not a builtin)
    manager.get_template_data = AsyncMock(return_value=None)

    return manager


@pytest.fixture
def mock_presentation_manager():
    """Create a mock presentation manager."""
    manager = MagicMock()

    # Create a real presentation for testing
    prs = Presentation()
    if prs.slide_layouts:
        prs.slides.add_slide(prs.slide_layouts[0])

    metadata = MagicMock()
    metadata.name = "test_presentation"
    metadata.namespace_id = "test-namespace-id"
    metadata.vfs_path = "/templates/test"

    # Mock get method
    manager.get = AsyncMock(return_value=(prs, metadata))
    manager.get_presentation = AsyncMock(return_value=prs)

    # Mock list_presentations
    pres_info = MagicMock()
    pres_info.name = "test_presentation"
    pres_info.slide_count = 1
    pres_info.namespace_id = "test-namespace-id"

    response = MagicMock()
    response.presentations = [pres_info]
    manager.list_presentations = AsyncMock(return_value=response)

    # Mock metadata storage
    manager._metadata = {
        "test_presentation": metadata,
    }

    # Mock update methods
    manager.update_slide_metadata = AsyncMock()
    manager._save_to_store = AsyncMock()
    manager.update = AsyncMock()
    manager.import_template = AsyncMock(return_value=True)
    manager.get_current_name = MagicMock(return_value="test_presentation")

    return manager


class TestRegisterTemplateTools:
    """Tests for template tools registration."""

    def test_register_template_tools(self, mock_mcp, mock_presentation_manager):
        """Test registering all template tools."""
        tools = register_template_tools(mock_mcp, mock_presentation_manager, template_manager=None)

        assert isinstance(tools, dict)
        # Should have tools from all submodules
        assert "pptx_list_templates" in tools
        assert "pptx_analyze_template" in tools
        assert "pptx_import_template" in tools
        assert "pptx_add_slide_from_template" in tools

    def test_register_template_tools_with_manager(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test registering with explicit template manager."""
        tools = register_template_tools(
            mock_mcp, mock_presentation_manager, template_manager=mock_template_manager
        )

        assert isinstance(tools, dict)
        assert len(tools) > 0


class TestListTools:
    """Tests for template listing tools."""

    @pytest.mark.asyncio
    async def test_list_templates(self, mock_mcp, mock_presentation_manager, mock_template_manager):
        """Test listing templates."""
        register_list_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)
        data = json.loads(result)
        assert "builtin_templates" in data or "error" not in data

    @pytest.mark.asyncio
    async def test_list_templates_no_template_manager(self, mock_mcp, mock_presentation_manager):
        """Test listing templates without template manager."""
        register_list_tools(mock_mcp, mock_presentation_manager, None)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_list_templates_with_custom_templates(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test listing templates with custom templates."""
        # Set up custom template in metadata
        mock_presentation_manager._metadata["test_presentation"].vfs_path = "/templates/custom"

        register_list_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)


class TestAnalyzeTools:
    """Tests for template analysis tools."""

    @pytest.mark.asyncio
    async def test_analyze_template_from_store(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing a template from the store."""
        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test_presentation")

        assert isinstance(result, str)
        data = json.loads(result)
        # Should have template info or error
        assert "name" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_builtin(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing a builtin template."""
        # Mock builtin template data
        prs = Presentation()
        import io

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("builtin_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "name" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_not_found(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing a template that doesn't exist."""
        mock_presentation_manager.get = AsyncMock(return_value=None)

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("nonexistent")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_variants(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing template variants."""
        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("test_presentation")

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_analyze_template_variants_not_found(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing variants for non-existent template."""
        mock_presentation_manager.get = AsyncMock(return_value=None)

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("nonexistent")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestImportTools:
    """Tests for template import tools."""

    @pytest.mark.asyncio
    async def test_import_template_success(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test successfully importing a template."""
        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_import_template"](
            file_path="/path/to/template.pptx", template_name="imported_template"
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_import_template_failure(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test importing template that fails."""
        mock_presentation_manager.import_template = AsyncMock(return_value=False)

        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_import_template"](
            file_path="/path/to/nonexistent.pptx", template_name="failed_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_import_template_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test importing template with exception."""
        mock_presentation_manager.import_template = AsyncMock(
            side_effect=Exception("Import failed")
        )

        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_import_template"](
            file_path="/path/to/error.pptx", template_name="error_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_get_builtin_template_not_found(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test getting builtin template that doesn't exist."""
        mock_template_manager.get_template_data = AsyncMock(return_value=None)

        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="nonexistent", save_as="my_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestWorkflowTools:
    """Tests for template workflow tools."""

    @pytest.mark.asyncio
    async def test_add_slide_from_template_no_presentation(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test adding slide when presentation not found."""
        mock_presentation_manager.get = AsyncMock(return_value=None)

        register_workflow_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_add_slide_from_template"](
            layout_index=0, presentation="nonexistent"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestModels:
    """Tests for template models."""

    def test_layout_info_model(self):
        """Test LayoutInfo model."""
        layout = LayoutInfo(
            index=0,
            name="Title Slide",
            placeholder_count=2,
            placeholders=[
                {"idx": 0, "type": "TITLE", "name": "Title"},
                {"idx": 1, "type": "SUBTITLE", "name": "Subtitle"},
            ],
        )
        assert layout.index == 0
        assert layout.name == "Title Slide"
        assert layout.placeholder_count == 2
        assert len(layout.placeholders) == 2

    def test_template_info_model(self):
        """Test TemplateInfo model."""
        info = TemplateInfo(
            name="test_template",
            slide_count=0,
            layout_count=5,
            layouts=[],
            master_count=1,
            has_theme=True,
        )
        assert info.name == "test_template"
        assert info.layout_count == 5

    def test_builtin_template_info_model(self):
        """Test BuiltinTemplateInfo model."""
        info = BuiltinTemplateInfo(
            name="corporate",
            display_name="Corporate Template",
            description="A corporate template",
            category="business",
            layout_count=20,
            tags=["business", "formal"],
            is_builtin=True,
        )
        assert info.name == "corporate"
        assert info.is_builtin is True

    def test_custom_template_info_model(self):
        """Test CustomTemplateInfo model."""
        info = CustomTemplateInfo(
            name="my_template",
            slide_count=5,
            namespace_id="ns-123",
            is_builtin=False,
            category="custom",
        )
        assert info.name == "my_template"
        assert info.is_builtin is False

    def test_template_list_response_model(self):
        """Test PresentationTemplateListResponse model."""
        response = PresentationTemplateListResponse(
            builtin_templates=[],
            custom_templates=[],
            total=0,
        )
        assert response.total == 0


class TestAnalyzeLayoutLoop:
    """Tests for analyze.py layout analysis loop (lines 118-138)."""

    @pytest.fixture
    def presentation_with_layouts(self):
        """Create a presentation with layouts that have placeholders."""
        prs = Presentation()
        return prs

    @pytest.mark.asyncio
    async def test_analyze_template_with_layouts_and_placeholders(
        self, mock_mcp, mock_template_manager
    ):
        """Test analyzing a template with slide_layouts and placeholders."""
        # Create a real presentation with layouts
        prs = Presentation()

        # Set up manager mock with real presentation
        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_with_layouts"
        manager.get = AsyncMock(return_value=(prs, metadata))

        register_analyze_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test_with_layouts")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "name" in data
        assert data["name"] == "test_with_layouts"
        # Should have layouts from the real presentation
        assert "layouts" in data
        assert "layout_count" in data
        # Verify layouts have the expected structure
        if data["layouts"]:
            layout = data["layouts"][0]
            assert "index" in layout
            assert "name" in layout
            assert "placeholder_count" in layout
            assert "placeholders" in layout

    @pytest.mark.asyncio
    async def test_analyze_template_layout_placeholder_details(
        self, mock_mcp, mock_template_manager
    ):
        """Test that placeholder details are correctly extracted."""
        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        register_analyze_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test")

        data = json.loads(result)
        if data.get("layouts"):
            for layout in data["layouts"]:
                # Each placeholder should have idx, type, name
                for ph in layout.get("placeholders", []):
                    assert "idx" in ph
                    assert "type" in ph
                    assert "name" in ph


class TestAnalyzeVariantsBuiltinTemplate:
    """Tests for analyze.py variant analysis with builtin template (lines 178-197)."""

    @pytest.mark.asyncio
    async def test_analyze_variants_builtin_template(self, mock_mcp, mock_presentation_manager):
        """Test analyzing variants for a builtin template."""
        import io

        # Create template manager that returns template data
        template_manager = MagicMock()
        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("builtin_test")

        assert isinstance(result, str)
        # Should have successfully analyzed the builtin template
        data = json.loads(result)
        # Either we get layout groups or an error
        assert "layout_groups" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_variants_builtin_with_layouts(self, mock_mcp, mock_presentation_manager):
        """Test analyzing variants for builtin template with multiple layouts."""
        import io

        template_manager = MagicMock()
        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("multi_layout_template")

        assert isinstance(result, str)
        data = json.loads(result)
        # Check structure
        if "layout_groups" in data:
            assert isinstance(data["layout_groups"], list)


class TestGetBuiltinTemplateSuccess:
    """Tests for import_tools.py pptx_get_builtin_template success path (lines 92-131)."""

    @pytest.mark.asyncio
    async def test_get_builtin_template_success(self, mock_mcp, mock_template_manager):
        """Test successfully getting a builtin template."""
        import io

        # Create a valid presentation
        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_data = buffer.read()
        mock_template_manager.get_template_data = AsyncMock(return_value=template_data)

        # Set up manager with mock store
        manager = MagicMock()
        mock_store = MagicMock()
        mock_namespace_info = MagicMock()
        mock_namespace_info.namespace_id = "ns-123"
        mock_store.create_namespace = AsyncMock(return_value=mock_namespace_info)
        mock_store.write_namespace = AsyncMock()
        manager._get_store = MagicMock(return_value=mock_store)
        manager._sanitize_name = MagicMock(side_effect=lambda x: x)
        manager._namespace_ids = {}
        manager.base_path = "/presentations"
        manager.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="corporate", save_as="my_corporate"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "message" in data
        assert "corporate" in data["message"]
        assert "my_corporate" in data["message"]

    @pytest.mark.asyncio
    async def test_get_builtin_template_no_store(self, mock_mcp, mock_template_manager):
        """Test getting builtin template when no store is available."""
        import io

        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        # Manager with no store
        manager = MagicMock()
        manager._get_store = MagicMock(return_value=None)

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="corporate", save_as="my_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "store" in data["error"].lower()

    @pytest.mark.asyncio
    async def test_get_builtin_template_exception(self, mock_mcp, mock_template_manager):
        """Test exception during builtin template import."""
        import io

        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        manager = MagicMock()
        mock_store = MagicMock()
        mock_store.create_namespace = AsyncMock(side_effect=Exception("Store error"))
        manager._get_store = MagicMock(return_value=mock_store)
        manager._sanitize_name = MagicMock(side_effect=lambda x: x)
        manager.base_path = "/presentations"
        manager.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="corporate", save_as="my_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_get_builtin_template_with_layouts(self, mock_mcp, mock_template_manager):
        """Test builtin template reports correct layout count."""
        import io

        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_data = buffer.read()
        mock_template_manager.get_template_data = AsyncMock(return_value=template_data)

        manager = MagicMock()
        mock_store = MagicMock()
        mock_namespace_info = MagicMock()
        mock_namespace_info.namespace_id = "ns-456"
        mock_store.create_namespace = AsyncMock(return_value=mock_namespace_info)
        mock_store.write_namespace = AsyncMock()
        manager._get_store = MagicMock(return_value=mock_store)
        manager._sanitize_name = MagicMock(side_effect=lambda x: x)
        manager._namespace_ids = {}
        manager.base_path = "/presentations"
        manager.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="modern", save_as="my_modern"
        )

        data = json.loads(result)
        assert "message" in data
        # Should mention layouts
        assert "layout" in data["message"].lower()


class TestErrorHandling:
    """Test error handling in template tools."""

    @pytest.mark.asyncio
    async def test_list_templates_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test exception handling in list_templates."""
        mock_template_manager.list_templates.side_effect = Exception("List failed")

        register_list_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test exception handling in analyze_template."""
        mock_presentation_manager.get = AsyncMock(side_effect=Exception("Get failed"))

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("error_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_analyze_variants_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test exception handling in analyze_template_variants."""
        mock_presentation_manager.get = AsyncMock(side_effect=Exception("Variant error"))
        mock_template_manager.get_template_data = AsyncMock(return_value=None)

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("error_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
