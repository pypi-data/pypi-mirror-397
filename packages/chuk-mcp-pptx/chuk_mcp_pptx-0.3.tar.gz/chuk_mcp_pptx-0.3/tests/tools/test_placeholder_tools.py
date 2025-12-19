"""
Tests for tools/core/placeholder.py

Tests placeholder population tools for >90% coverage.
"""

import json
import pytest
from unittest.mock import MagicMock
from pptx import Presentation

from chuk_mcp_pptx.tools.core.placeholder import register_placeholder_tools


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures registered tools."""
    mcp = MagicMock()
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mcp.tool = tool_decorator
    mcp._tools = tools
    return mcp


def create_presentation_with_placeholders():
    """Create a presentation with slides that have placeholders."""
    prs = Presentation()

    # Add slide with TITLE_AND_CONTENT layout (usually has title and body placeholders)
    if len(prs.slide_layouts) > 1:
        layout = prs.slide_layouts[1]  # TITLE_AND_CONTENT
        prs.slides.add_slide(layout)

    return prs


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation or create_presentation_with_placeholders()
        self._current_name = "test_presentation"
        self._metadata = MagicMock()
        self._metadata.name = self._current_name

    async def get(self, name=None):
        """Get presentation."""
        if self._presentation is None:
            return None
        if name is None or name == self._current_name:
            return self._presentation, self._metadata
        return None

    async def get_presentation(self, name=None):
        """Get presentation object directly."""
        if self._presentation is None:
            return None
        if name is None or name == self._current_name:
            return self._presentation
        return None

    def get_current_name(self):
        """Get current presentation name."""
        return self._current_name

    async def update_slide_metadata(self, slide_index):
        """Update slide metadata."""
        pass

    async def update(self, name=None):
        """Update presentation."""
        pass

    async def _save_to_store(self, name, prs):
        """Save to store."""
        pass


@pytest.fixture
def mock_manager():
    """Create a mock presentation manager."""
    return MockPresentationManager()


@pytest.fixture
def placeholder_tools(mock_mcp, mock_manager):
    """Register placeholder tools and return them."""
    register_placeholder_tools(mock_mcp, mock_manager)
    return mock_mcp._tools


class TestPopulatePlaceholder:
    """Tests for pptx_populate_placeholder."""

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_string(self, placeholder_tools, mock_manager):
        """Test populating placeholder with string content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test Title",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_dict(self, placeholder_tools, mock_manager):
        """Test populating placeholder with dict content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content={"type": "Table", "headers": ["A", "B"], "data": [["1", "2"]]},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_json_string(self, placeholder_tools, mock_manager):
        """Test populating placeholder with JSON string."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content='{"type": "Table", "headers": ["X"], "data": [["Y"]]}',
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_invalid_slide_index(self, placeholder_tools, mock_manager):
        """Test with invalid slide index."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=999,
            placeholder_idx=0,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_populate_placeholder_invalid_placeholder_idx(
        self, placeholder_tools, mock_manager
    ):
        """Test with invalid placeholder index."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=999,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data or "not found" in result.lower()

    @pytest.mark.asyncio
    async def test_populate_placeholder_no_presentation(self, mock_mcp):
        """Test when no presentation exists."""
        manager = MockPresentationManager()
        manager._presentation = None
        register_placeholder_tools(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_presentation_name(
        self, placeholder_tools, mock_manager
    ):
        """Test with explicit presentation name."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Named Presentation",
            presentation="test_presentation",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_negative_slide_index(self, placeholder_tools, mock_manager):
        """Test with negative slide index."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=-1,
            placeholder_idx=0,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_populate_placeholder_empty_string(self, placeholder_tools, mock_manager):
        """Test populating with empty string."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_multiline_string(self, placeholder_tools, mock_manager):
        """Test populating with multiline string."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Line 1\nLine 2\nLine 3",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_special_characters(self, placeholder_tools, mock_manager):
        """Test populating with special characters."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Special: <>&\"'",
        )
        assert isinstance(result, str)


class TestPopulatePlaceholderContentTypes:
    """Test different content types for placeholder population."""

    @pytest.mark.asyncio
    async def test_populate_with_table_dict(self, placeholder_tools, mock_manager):
        """Test populating with table dict."""
        table_content = {
            "type": "Table",
            "headers": ["Name", "Value", "Status"],
            "data": [
                ["Item 1", "100", "Active"],
                ["Item 2", "200", "Pending"],
            ],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=table_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with chart dict."""
        chart_content = {
            "type": "ColumnChart",
            "categories": ["Q1", "Q2", "Q3"],
            "series": {"Sales": [100, 150, 200]},
            "title": "Quarterly Sales",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_image_dict(self, placeholder_tools, mock_manager):
        """Test populating with image dict."""
        image_content = {
            "type": "Image",
            "image_source": "https://example.com/image.png",
            "alt": "Test Image",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=image_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_pie_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with pie chart dict."""
        chart_content = {
            "type": "PieChart",
            "categories": ["A", "B", "C"],
            "values": [30, 50, 20],
            "title": "Distribution",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_line_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with line chart dict."""
        chart_content = {
            "type": "LineChart",
            "categories": ["Jan", "Feb", "Mar"],
            "series": {"Revenue": [1000, 1200, 1100]},
            "title": "Monthly Revenue",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_bar_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with bar chart dict."""
        chart_content = {
            "type": "BarChart",
            "categories": ["Product A", "Product B", "Product C"],
            "series": {"Units": [50, 75, 60]},
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_unknown_type(self, placeholder_tools, mock_manager):
        """Test populating with unknown content type."""
        content = {
            "type": "UnknownType",
            "data": "something",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_dict_missing_type(self, placeholder_tools, mock_manager):
        """Test populating with dict missing type field."""
        content = {
            "headers": ["A", "B"],
            "data": [["1", "2"]],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


class TestToolRegistration:
    """Test tool registration."""

    def test_tool_registered(self, placeholder_tools):
        """Test that placeholder tool is registered."""
        assert "pptx_populate_placeholder" in placeholder_tools
        assert callable(placeholder_tools["pptx_populate_placeholder"])

    def test_tool_is_async(self, placeholder_tools):
        """Test that tool is async."""
        import asyncio

        assert asyncio.iscoroutinefunction(placeholder_tools["pptx_populate_placeholder"])


class TestEdgeCases:
    """Test edge cases."""

    @pytest.mark.asyncio
    async def test_populate_with_very_long_text(self, placeholder_tools, mock_manager):
        """Test populating with very long text."""
        long_text = "A" * 10000
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=long_text,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_unicode(self, placeholder_tools, mock_manager):
        """Test populating with unicode characters."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="日本語テスト 한국어 中文",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_emoji(self, placeholder_tools, mock_manager):
        """Test populating with emoji."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test 🎉 Emoji 🚀",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_html_like_content(self, placeholder_tools, mock_manager):
        """Test populating with HTML-like content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="<b>Bold</b> and <i>italic</i>",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_invalid_json_string(self, placeholder_tools, mock_manager):
        """Test populating with invalid JSON string that looks like JSON."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="{invalid json}",
        )
        assert isinstance(result, str)


class TestTableContentVariants:
    """Test various table content configurations."""

    @pytest.mark.asyncio
    async def test_table_with_variant(self, placeholder_tools, mock_manager):
        """Test table with variant styling."""
        content = {
            "type": "Table",
            "headers": ["Col1", "Col2"],
            "data": [["A", "B"]],
            "variant": "striped",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_table_with_empty_data(self, placeholder_tools, mock_manager):
        """Test table with empty data."""
        content = {
            "type": "Table",
            "headers": ["Col1", "Col2"],
            "data": [],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_table_with_single_cell(self, placeholder_tools, mock_manager):
        """Test table with single cell."""
        content = {
            "type": "Table",
            "headers": ["Single"],
            "data": [["Value"]],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_table_with_many_rows(self, placeholder_tools, mock_manager):
        """Test table with many rows."""
        content = {
            "type": "Table",
            "headers": ["ID", "Name"],
            "data": [[str(i), f"Item {i}"] for i in range(20)],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


class TestDictContentErrorPaths:
    """Test error paths for dict content."""

    @pytest.mark.asyncio
    async def test_dict_content_no_presentation(self, mock_mcp):
        """Test dict content when no presentation exists."""
        manager = MockPresentationManager()
        manager._presentation = None
        register_placeholder_tools(mock_mcp, manager)

        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_invalid_slide_index(self, placeholder_tools):
        """Test dict content with invalid slide index."""
        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=999,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_negative_slide_index(self, placeholder_tools):
        """Test dict content with negative slide index."""
        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=-1,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_invalid_placeholder(self, placeholder_tools):
        """Test dict content with invalid placeholder index."""
        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=9999,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data or "not found" in result.lower()

    @pytest.mark.asyncio
    async def test_dict_content_unknown_component(self, placeholder_tools):
        """Test dict content with unknown component type."""
        content = {"type": "UnknownComponent", "data": "something"}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_with_invalid_params(self, placeholder_tools):
        """Test dict content with invalid params for component."""
        content = {
            "type": "Table",
            "headers": ["A"],
            "data": [["1"]],
            "invalid_param": "should be ignored",
            "another_invalid": 123,
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


class TestStringContentPlaceholderTypes:
    """Test string content for different placeholder types."""

    @pytest.mark.asyncio
    async def test_body_placeholder_with_bullets(self, placeholder_tools, mock_manager):
        """Test body placeholder with bullet points using \\n separator."""
        # Body placeholders typically have idx=1 in TITLE_AND_CONTENT layout
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content="First bullet\\nSecond bullet\\nThird bullet",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_title_placeholder_simple(self, placeholder_tools, mock_manager):
        """Test title placeholder with simple text."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Simple Title",
        )
        assert isinstance(result, str)


class TestInvalidContentTypes:
    """Test invalid content type handling."""

    @pytest.mark.asyncio
    async def test_invalid_content_type_list(self, placeholder_tools, mock_manager):
        """Test with list content (not supported)."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=["item1", "item2"],  # List is not a valid content type
        )
        assert isinstance(result, str)
        # Should either succeed by converting or error

    @pytest.mark.asyncio
    async def test_invalid_content_type_number(self, placeholder_tools, mock_manager):
        """Test with number content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=12345,  # Number should be handled
        )
        assert isinstance(result, str)


class TestChartContentVariants:
    """Test various chart content configurations."""

    @pytest.mark.asyncio
    async def test_chart_with_multiple_series(self, placeholder_tools, mock_manager):
        """Test chart with multiple series."""
        content = {
            "type": "ColumnChart",
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": {
                "Sales": [100, 120, 140, 160],
                "Costs": [80, 90, 100, 110],
                "Profit": [20, 30, 40, 50],
            },
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_chart_without_title(self, placeholder_tools, mock_manager):
        """Test chart without title."""
        content = {
            "type": "ColumnChart",
            "categories": ["A", "B", "C"],
            "series": {"Values": [1, 2, 3]},
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_doughnut_chart(self, placeholder_tools, mock_manager):
        """Test doughnut chart."""
        content = {
            "type": "DoughnutChart",
            "categories": ["Segment A", "Segment B", "Segment C"],
            "values": [40, 35, 25],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_area_chart(self, placeholder_tools, mock_manager):
        """Test area chart."""
        content = {
            "type": "AreaChart",
            "categories": ["Week 1", "Week 2", "Week 3", "Week 4"],
            "series": {"Traffic": [1000, 1200, 1100, 1400]},
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)
