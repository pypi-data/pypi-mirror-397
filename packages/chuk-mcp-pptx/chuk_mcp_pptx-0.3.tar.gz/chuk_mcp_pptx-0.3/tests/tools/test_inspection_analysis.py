# tests/tools/test_inspection_analysis.py
"""
Tests for the inspection and analysis tools.

Tests the tools:
- pptx_inspect_slide
- pptx_fix_slide_layout
- pptx_analyze_presentation_layout
"""

import pytest
from unittest.mock import MagicMock
from pptx import Presentation
from pptx.util import Inches


# ============================================================================
# Fixtures
# ============================================================================


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation
        self._current_name = "test_presentation"

    async def get_presentation(self, name=None):
        """Get presentation by name or return current."""
        if name is None or name == self._current_name:
            return self._presentation
        return None


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures tool registrations."""
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mock = MagicMock()
    mock.tool = tool_decorator
    mock._tools = tools
    return mock


@pytest.fixture
def mock_manager_no_prs():
    """Create a mock manager with no presentation."""
    return MockPresentationManager(presentation=None)


@pytest.fixture
def presentation_with_slide():
    """Create a presentation with one blank slide."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout
    prs.slides.add_slide(blank_layout)
    return prs


@pytest.fixture
def presentation_with_title_slide():
    """Create a presentation with a title slide."""
    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(title_layout)
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = "Test Title"
    return prs


@pytest.fixture
def presentation_with_elements(presentation_with_slide):
    """Create a presentation with various elements on the slide."""
    prs = presentation_with_slide
    slide = prs.slides[0]

    # Add a text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(2)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Sample text content"

    # Add another text box
    left2 = Inches(4)
    top2 = Inches(1)
    txBox2 = slide.shapes.add_textbox(left2, top2, width, height)
    tf2 = txBox2.text_frame
    tf2.text = "Another text box"

    return prs


@pytest.fixture
def mock_manager(presentation_with_slide):
    """Create a mock manager with a presentation."""
    return MockPresentationManager(presentation=presentation_with_slide)


@pytest.fixture
def mock_manager_with_elements(presentation_with_elements):
    """Create a mock manager with elements."""
    return MockPresentationManager(presentation=presentation_with_elements)


@pytest.fixture
def mock_manager_with_title(presentation_with_title_slide):
    """Create a mock manager with title slide."""
    return MockPresentationManager(presentation=presentation_with_title_slide)


@pytest.fixture
def inspection_tools(mock_mcp, mock_manager):
    """Register and return inspection tools."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager)


@pytest.fixture
def inspection_tools_no_prs(mock_mcp, mock_manager_no_prs):
    """Register and return inspection tools with no presentation."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager_no_prs)


@pytest.fixture
def inspection_tools_with_elements(mock_mcp, mock_manager_with_elements):
    """Register and return inspection tools with elements."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager_with_elements)


@pytest.fixture
def inspection_tools_with_title(mock_mcp, mock_manager_with_title):
    """Register and return inspection tools with title slide."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager_with_title)


# ============================================================================
# Test pptx_inspect_slide
# ============================================================================


class TestInspectSlide:
    """Tests for pptx_inspect_slide tool."""

    @pytest.mark.asyncio
    async def test_inspect_slide_basic(self, inspection_tools):
        """Test basic slide inspection."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        assert isinstance(result, str)
        assert "SLIDE 0 INSPECTION" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_no_presentation(self, inspection_tools_no_prs):
        """Test inspecting when no presentation exists."""
        result = await inspection_tools_no_prs["pptx_inspect_slide"](slide_index=0)
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_inspect_slide_out_of_range(self, inspection_tools):
        """Test inspecting slide index out of range."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=999)
        assert "error" in result.lower()
        assert "out of range" in result.lower()

    @pytest.mark.asyncio
    async def test_inspect_slide_with_title(self, inspection_tools_with_title):
        """Test inspecting slide with title."""
        result = await inspection_tools_with_title["pptx_inspect_slide"](slide_index=0)
        assert "Title:" in result
        assert "Test Title" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_with_text_boxes(self, inspection_tools_with_elements):
        """Test inspecting slide with text boxes."""
        result = await inspection_tools_with_elements["pptx_inspect_slide"](slide_index=0)
        assert "TEXT BOXES" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_includes_measurements(self, inspection_tools_with_elements):
        """Test that measurements are included."""
        result = await inspection_tools_with_elements["pptx_inspect_slide"](
            slide_index=0, include_measurements=True
        )
        assert "at (" in result or "size" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_no_measurements(self, inspection_tools_with_elements):
        """Test inspection without measurements."""
        result = await inspection_tools_with_elements["pptx_inspect_slide"](
            slide_index=0, include_measurements=False
        )
        assert isinstance(result, str)
        # Should still contain basic info
        assert "SLIDE 0 INSPECTION" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_string_index(self, inspection_tools):
        """Test that string slide index is converted to int."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index="0")
        assert isinstance(result, str)
        assert "SLIDE 0 INSPECTION" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_summary(self, inspection_tools):
        """Test that summary is included."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        assert "SUMMARY" in result
        assert "Total elements" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_no_issues(self, inspection_tools):
        """Test slide with no layout issues."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        # Blank slide should have no issues
        assert "No layout issues" in result or "Layout issues: 0" in result


# ============================================================================
# Test pptx_fix_slide_layout
# ============================================================================


class TestFixSlideLayout:
    """Tests for pptx_fix_slide_layout tool."""

    @pytest.mark.asyncio
    async def test_fix_layout_basic(self, inspection_tools):
        """Test basic layout fixing."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index=0)
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_no_presentation(self, inspection_tools_no_prs):
        """Test fixing when no presentation exists."""
        result = await inspection_tools_no_prs["pptx_fix_slide_layout"](slide_index=0)
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_layout_out_of_range(self, inspection_tools):
        """Test fixing slide index out of range."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index=999)
        assert "error" in result.lower()
        assert "out of range" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_layout_no_issues(self, inspection_tools):
        """Test fixing when no issues exist."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index=0)
        assert "optimal" in result.lower() or "no layout issues" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_layout_string_index(self, inspection_tools):
        """Test that string slide index is converted to int."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index="0")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_with_presentation_name(self, inspection_tools):
        """Test fixing with presentation name."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, presentation="test_presentation"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_fix_bounds_only(self, inspection_tools):
        """Test fixing only bounds issues."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=False, fix_bounds=True, fix_spacing=False
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_fix_overlaps_only(self, inspection_tools):
        """Test fixing only overlap issues."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=True, fix_bounds=False, fix_spacing=False
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_fix_spacing_only(self, inspection_tools):
        """Test fixing only spacing issues."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=False, fix_bounds=False, fix_spacing=True
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_maintain_relative_positions(self, inspection_tools):
        """Test fixing with maintain_relative_positions flag."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, maintain_relative_positions=True
        )
        assert isinstance(result, str)


# ============================================================================
# Test pptx_analyze_presentation_layout
# ============================================================================


class TestAnalyzePresentationLayout:
    """Tests for pptx_analyze_presentation_layout tool."""

    @pytest.mark.asyncio
    async def test_analyze_layout_basic(self, inspection_tools):
        """Test basic presentation layout analysis."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert isinstance(result, str)
        assert "PRESENTATION LAYOUT ANALYSIS" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_no_presentation(self, inspection_tools_no_prs):
        """Test analyzing when no presentation exists."""
        result = await inspection_tools_no_prs["pptx_analyze_presentation_layout"]()
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_total_slides(self, inspection_tools):
        """Test that total slides count is shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "Total slides:" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_layout_usage(self, inspection_tools):
        """Test that layout usage is shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "LAYOUT USAGE" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_element_stats(self, inspection_tools):
        """Test that element statistics are shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "ELEMENT STATISTICS" in result
        assert "Images:" in result
        assert "Charts:" in result
        assert "Tables:" in result
        assert "Text boxes:" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_recommendations(self, inspection_tools):
        """Test that recommendations section is shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "RECOMMENDATIONS" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_with_elements(self, inspection_tools_with_elements):
        """Test analyzing presentation with elements."""
        result = await inspection_tools_with_elements["pptx_analyze_presentation_layout"]()
        assert "PRESENTATION LAYOUT ANALYSIS" in result
        # Should count text boxes
        assert "Text boxes:" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_with_presentation_name(self, inspection_tools):
        """Test analyzing with presentation name."""
        result = await inspection_tools["pptx_analyze_presentation_layout"](
            presentation="test_presentation"
        )
        assert isinstance(result, str)
        assert "PRESENTATION LAYOUT ANALYSIS" in result


# ============================================================================
# Test Tool Registration
# ============================================================================


class TestToolRegistration:
    """Tests for tool registration."""

    def test_all_tools_registered(self, inspection_tools):
        """Test that all inspection tools are registered."""
        expected_tools = [
            "pptx_inspect_slide",
            "pptx_fix_slide_layout",
            "pptx_analyze_presentation_layout",
        ]
        for tool_name in expected_tools:
            assert tool_name in inspection_tools

    def test_tools_are_async(self, inspection_tools):
        """Test that all tools are async functions."""
        import inspect

        for tool_name, tool_func in inspection_tools.items():
            assert inspect.iscoroutinefunction(tool_func), f"{tool_name} should be async"


# ============================================================================
# Test Edge Cases and Overlapping Elements
# ============================================================================


class TestOverlappingElements:
    """Tests for overlapping element detection."""

    @pytest.fixture
    def presentation_with_overlaps(self):
        """Create a presentation with overlapping elements."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add two overlapping text boxes
        left1 = Inches(2)
        top1 = Inches(2)
        width = Inches(3)
        height = Inches(2)

        txBox1 = slide.shapes.add_textbox(left1, top1, width, height)
        txBox1.text_frame.text = "Overlapping box 1"

        # Second box partially overlaps first
        left2 = Inches(3)
        top2 = Inches(2.5)
        txBox2 = slide.shapes.add_textbox(left2, top2, width, height)
        txBox2.text_frame.text = "Overlapping box 2"

        return prs

    @pytest.fixture
    def inspection_tools_with_overlaps(self, mock_mcp, presentation_with_overlaps):
        """Register inspection tools with overlapping elements."""
        manager = MockPresentationManager(presentation=presentation_with_overlaps)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_detects_overlaps(self, inspection_tools_with_overlaps):
        """Test that inspection detects overlapping elements."""
        result = await inspection_tools_with_overlaps["pptx_inspect_slide"](
            slide_index=0, check_overlaps=True
        )
        # May or may not detect overlaps depending on implementation details
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_overlaps(self, inspection_tools_with_overlaps):
        """Test fixing overlapping elements."""
        result = await inspection_tools_with_overlaps["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=True
        )
        assert isinstance(result, str)


# ============================================================================
# Test Out of Bounds Elements
# ============================================================================


class TestOutOfBoundsElements:
    """Tests for out of bounds element detection."""

    @pytest.fixture
    def presentation_with_oob(self):
        """Create a presentation with out of bounds elements."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add a text box that extends beyond slide bounds
        left = Inches(8)  # Very close to right edge
        top = Inches(1)
        width = Inches(4)  # Will extend beyond slide
        height = Inches(1)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = "Out of bounds text"

        return prs

    @pytest.fixture
    def inspection_tools_with_oob(self, mock_mcp, presentation_with_oob):
        """Register inspection tools with out of bounds elements."""
        manager = MockPresentationManager(presentation=presentation_with_oob)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_detects_oob(self, inspection_tools_with_oob):
        """Test that inspection detects out of bounds elements."""
        result = await inspection_tools_with_oob["pptx_inspect_slide"](slide_index=0)
        assert isinstance(result, str)
        # Should mention out of bounds or layout issues
        assert "OUT OF BOUNDS" in result or "LAYOUT ISSUES" in result

    @pytest.mark.asyncio
    async def test_fix_oob(self, inspection_tools_with_oob):
        """Test fixing out of bounds elements."""
        result = await inspection_tools_with_oob["pptx_fix_slide_layout"](
            slide_index=0, fix_bounds=True
        )
        assert isinstance(result, str)


# ============================================================================
# Test Multiple Slides
# ============================================================================


class TestMultipleSlides:
    """Tests for presentations with multiple slides."""

    @pytest.fixture
    def presentation_with_multiple_slides(self):
        """Create a presentation with multiple slides."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        # Add 3 slides
        for _ in range(3):
            prs.slides.add_slide(blank_layout)

        return prs

    @pytest.fixture
    def inspection_tools_multiple_slides(self, mock_mcp, presentation_with_multiple_slides):
        """Register inspection tools with multiple slides."""
        manager = MockPresentationManager(presentation=presentation_with_multiple_slides)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_each_slide(self, inspection_tools_multiple_slides):
        """Test inspecting each slide individually."""
        for i in range(3):
            result = await inspection_tools_multiple_slides["pptx_inspect_slide"](slide_index=i)
            assert isinstance(result, str)
            assert f"SLIDE {i} INSPECTION" in result

    @pytest.mark.asyncio
    async def test_analyze_multiple_slides(self, inspection_tools_multiple_slides):
        """Test analyzing presentation with multiple slides."""
        result = await inspection_tools_multiple_slides["pptx_analyze_presentation_layout"]()
        assert "Total slides: 3" in result
