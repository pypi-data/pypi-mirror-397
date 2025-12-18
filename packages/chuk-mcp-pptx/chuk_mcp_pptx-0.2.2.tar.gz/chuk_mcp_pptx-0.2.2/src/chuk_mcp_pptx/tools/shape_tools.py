# src/chuk_mcp_pptx/tools/shape_tools.py
"""
Shape and SmartArt Tools for PowerPoint MCP Server

Provides async MCP tools for creating shapes, connectors, and SmartArt-like diagrams.
"""

from typing import List, Optional


def register_shape_tools(mcp, manager):
    """Register all shape and SmartArt tools with the MCP server.

    Note: pptx_add_shape is provided by component_tools.py as part of the
    comprehensive design system. This module provides specialized shape tools.
    """

    from ..components.core import Shape, Connector, ProcessFlow, CycleDiagram, HierarchyDiagram
    from ..layout.helpers import validate_position
    from ..components.code import CodeBlock
    from ..themes.theme_manager import ThemeManager

    # pptx_add_shape removed - now in component_tools.py as part of design system
    # Keeping specialized tools: pptx_add_arrow, pptx_add_smart_art, pptx_add_code_block

    @mcp.tool
    async def pptx_add_arrow(
        slide_index: int,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        connector_type: str = "straight",
        line_color: Optional[str] = "#000000",
        line_width: float = 2.0,
        arrow_start: bool = False,
        arrow_end: bool = True,
        presentation: Optional[str] = None,
    ) -> str:
        """
        Add an arrow or connector line to a slide.

        Creates connectors between points with optional arrowheads.

        Args:
            slide_index: Index of the slide to add arrow to
            start_x: Starting X position in inches
            start_y: Starting Y position in inches
            end_x: Ending X position in inches
            end_y: Ending Y position in inches
            connector_type: Type of connector ("straight", "elbow", "curved")
            line_color: Line color as hex string
            line_width: Line width in points
            arrow_start: Whether to add arrowhead at start
            arrow_end: Whether to add arrowhead at end
            presentation: Name of presentation

        Returns:
            Success message confirming arrow addition

        Example:
            await pptx_add_arrow(
                slide_index=1,
                start_x=2.0,
                start_y=2.0,
                end_x=5.0,
                end_y=3.0,
                connector_type="straight",
                arrow_end=True
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return "Error: No presentation found"

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            # Get theme for connector
            theme_manager = ThemeManager()
            theme_obj = theme_manager.get_default_theme()

            # Create connector using component
            connector_comp = Connector(
                start_x=start_x,
                start_y=start_y,
                end_x=end_x,
                end_y=end_y,
                connector_type=connector_type,
                line_color=line_color,  # Component handles hex colors
                line_width=line_width,
                arrow_start=arrow_start,
                arrow_end=arrow_end,
                theme=theme_obj,
            )
            connector_comp.render(slide)

            # Update in VFS
            await manager.update(presentation)

            return f"Added {connector_type} arrow to slide {slide_index}"

        except Exception as e:
            return f"Error adding arrow: {str(e)}"

    @mcp.tool
    async def pptx_add_smart_art(
        slide_index: int,
        art_type: str,
        items: List[str],
        title: Optional[str] = None,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
        color_scheme: str = "modern_blue",
        presentation: Optional[str] = None,
    ) -> str:
        """
        Add a SmartArt-style diagram to a slide.

        Creates professional diagrams like process flows, cycles, hierarchies, etc.

        Args:
            slide_index: Index of the slide to add SmartArt to
            art_type: Type of SmartArt diagram:
                - "process" - Sequential process flow
                - "cycle" - Circular/cyclical process
                - "hierarchy" - Organizational hierarchy
                - "list" - Bulleted list with shapes
                - "relationship" - Relationship diagram
                - "pyramid" - Pyramid hierarchy
            items: List of text items for the diagram
            title: Optional title for the diagram
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            color_scheme: Color scheme ("modern_blue", "corporate_gray", "warm_orange")
            presentation: Name of presentation

        Returns:
            Success message confirming SmartArt addition

        Example:
            await pptx_add_smart_art(
                slide_index=2,
                art_type="process",
                items=["Research", "Design", "Develop", "Test", "Deploy"],
                title="Development Process",
                color_scheme="modern_blue"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return "Error: No presentation found"

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            # Remove content placeholders that might interfere with SmartArt
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, "shape_type"):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format.type in [
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ]:
                                continue
                        # Mark content placeholders for removal
                        placeholders_to_remove.append(shape)

            # Remove the placeholders
            for placeholder in placeholders_to_remove:
                sp = placeholder._element
                sp.getparent().remove(sp)

            # Validate position
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )

            # Add title if provided
            if title:
                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_default_theme()

                title_comp = Shape(
                    shape_type="rectangle",
                    text=title,
                    fill_color="transparent",
                    line_color="transparent",
                    theme=theme_obj,
                )
                title_comp.render(slide, validated_left, validated_top - 0.5, validated_width, 0.4)
                # Adjust top position for the SmartArt
                validated_top += 0.2
                validated_height -= 0.7

            # Get theme for SmartArt
            theme_manager = ThemeManager()
            theme_obj = theme_manager.get_default_theme()

            # Map art_type to component
            art_components = {
                "process": ProcessFlow,
                "cycle": CycleDiagram,
                "hierarchy": HierarchyDiagram,
            }

            component_class = art_components.get(art_type)
            if not component_class:
                return f"Error: Unsupported art type '{art_type}'. Supported: {', '.join(art_components.keys())}"

            # Create SmartArt using component
            smart_art_comp = component_class(items=items, theme=theme_obj)
            smart_art_comp.render(
                slide, validated_left, validated_top, validated_width, validated_height
            )

            # Update in VFS
            await manager.update(presentation)

            # Report if position was adjusted
            position_note = ""
            if (
                validated_left != left
                or validated_top != top
                or validated_width != width
                or validated_height != height
            ):
                position_note = " (position adjusted)"

            return f"Added {art_type} SmartArt with {len(items)} items to slide {slide_index}{position_note}"

        except Exception as e:
            return f"Error adding SmartArt: {str(e)}"

    @mcp.tool
    async def pptx_add_code_block(
        slide_index: int,
        code: str,
        language: str = "python",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
        theme: str = "dark_modern",
        presentation: Optional[str] = None,
    ) -> str:
        """
        Add a code block to a slide with syntax highlighting appearance.

        Creates a formatted code block with monospace font and theme colors.

        Args:
            slide_index: Index of the slide to add code to (0-based)
            code: The code content to display
            language: Programming language (for label)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            theme: Theme to use for code block styling
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming code block addition

        Example:
            await pptx_add_code_block(
                slide_index=1,
                code="def hello_world():\\n    print('Hello, World!')",
                language="python",
                theme="dark_purple"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return "Error: No presentation found"

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            # Remove content placeholders
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, "shape_type"):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format.type in [
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ]:
                                continue
                        placeholders_to_remove.append(shape)

            for shape in placeholders_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)

            # Add code block using CodeBlock component
            theme_manager = ThemeManager()
            theme_obj = theme_manager.get_theme(theme) if theme else None

            code_component = CodeBlock(code=code, language=language, theme=theme_obj)
            code_component.render(slide, left=left, top=top, width=width, height=height)

            # Update in VFS
            await manager.update(presentation)

            return f"Added {language} code block to slide {slide_index}"

        except Exception as e:
            return f"Error adding code block: {str(e)}"

    # Theme tools have been moved to tools/theme_tools.py for better organization
    # Use: pptx_apply_theme, pptx_list_themes, pptx_create_custom_theme

    # Return tools for external access
    # Note: pptx_add_shape is provided by component_tools.py
    # Note: Theme tools are provided by theme_tools.py
    return {
        "pptx_add_arrow": pptx_add_arrow,
        "pptx_add_smart_art": pptx_add_smart_art,
        "pptx_add_code_block": pptx_add_code_block,
    }
