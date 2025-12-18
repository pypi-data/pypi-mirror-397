# src/chuk_mcp_pptx/tools/table_tools.py
"""
Table Tools for PowerPoint MCP Server

Provides async MCP tools for creating and managing tables in presentations.
Includes layout validation and boundary checking like charts and images.
"""

from ..models import ErrorResponse, SuccessResponse, ComponentResponse
from ..constants import (
    ErrorMessages,
)

from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def register_table_tools(mcp, manager):
    """Register all table-related tools with the MCP server."""

    from ..layout.helpers import validate_position, get_safe_content_area
    from ..components.core import Table
    from ..themes.theme_manager import ThemeManager

    @mcp.tool
    async def pptx_add_data_table(
        slide_index: int,
        headers: list[str],
        data: list[list[str]],
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.0,
        style: str = "medium",
        presentation: str | None = None,
    ) -> str:
        """
        Add a formatted data table to a slide with layout validation.

        Creates a table with headers and data rows, with optional styling.
        Automatically validates position to ensure it fits within slide boundaries
        and doesn't overlap with titles or other placeholders.

        Args:
            slide_index: Index of the slide to add table to (0-based)
            headers: List of column header strings
            data: List of rows, where each row is a list of cell values
            left: Left position in inches (will be validated)
            top: Top position in inches (will be validated)
            width: Width in inches (will be validated)
            height: Height in inches (will be validated)
            style: Table style (light, medium, dark)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming table addition with final position

        Example:
            await pptx_add_data_table(
                slide_index=1,
                headers=["Product", "Q1", "Q2", "Q3", "Q4"],
                data=[
                    ["Laptops", "$100K", "$120K", "$110K", "$130K"],
                    ["Phones", "$80K", "$90K", "$95K", "$100K"]
                ],
                left=1.0,
                top=2.0,
                width=8.0,
                height=3.0,
                style="medium"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            slide = prs.slides[slide_index]

            # Get safe content area considering if there's a title
            safe_area = get_safe_content_area(has_title=bool(slide.shapes.title))

            # Validate and adjust position to fit within slide
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )

            # Further adjust if position is too close to title area
            if slide.shapes.title and validated_top < safe_area["top"]:
                validated_top = safe_area["top"]

            # Remove any overlapping placeholders (except title)
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, "shape_type"):
                    # Check if it's a placeholder (but not a title)
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, "placeholder_format"):
                            from pptx.enum.shapes import PP_PLACEHOLDER

                            if shape.placeholder_format.type in [
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ]:
                                continue

                        # Check if placeholder overlaps with table area
                        shape_left = shape.left.inches if hasattr(shape.left, "inches") else 0
                        shape_top = shape.top.inches if hasattr(shape.top, "inches") else 0
                        shape_right = shape_left + (
                            shape.width.inches if hasattr(shape.width, "inches") else 0
                        )
                        shape_bottom = shape_top + (
                            shape.height.inches if hasattr(shape.height, "inches") else 0
                        )

                        table_right = validated_left + validated_width
                        table_bottom = validated_top + validated_height

                        # Check for overlap
                        if not (
                            shape_right < validated_left
                            or shape_left > table_right
                            or shape_bottom < validated_top
                            or shape_top > table_bottom
                        ):
                            placeholders_to_remove.append(shape)

            # Remove overlapping placeholders
            for placeholder in placeholders_to_remove:
                slide.shapes._spTree.remove(placeholder.element)

            # Get current theme
            theme_manager = ThemeManager()
            theme = theme_manager.get_theme(metadata.theme) if metadata.theme else None

            # Map style to variant
            variant_map = {"light": "minimal", "medium": "default", "dark": "bordered"}
            variant = variant_map.get(style, "default")

            # Create and render table using new Table component
            table_comp = Table(headers=headers, data=data, variant=variant, size="md", theme=theme)

            table_comp.render(
                slide,
                left=validated_left,
                top=validated_top,
                width=validated_width,
                height=validated_height,
            )

            # Update in VFS
            await manager.update(presentation)

            # Report if position was adjusted
            position_note = ""
            if (
                validated_left != left
                or validated_top != top
                or validated_width != width
                or validated_height != height
            ):
                position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"

            message = f"Added {len(data)} row table to slide {slide_index}{position_note}"
            return ComponentResponse(
                presentation=metadata.name,
                slide_index=slide_index,
                component="table",
                message=message,
                variant=None,
            ).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=f"Error adding table: {str(e)}").model_dump_json()

    @mcp.tool
    async def pptx_add_comparison_table(
        slide_index: int,
        title: str,
        categories: list[str],
        option1_name: str,
        option1_values: list[str],
        option2_name: str,
        option2_values: list[str],
        option3_name: str | None = None,
        option3_values: list[str] | None = None,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.0,
        style: str = "light",
        presentation: str | None = None,
    ) -> str:
        """
        Add a comparison table for multiple options with layout validation.

        Creates a table comparing 2-3 options across multiple categories.
        Validates position to ensure proper fit within slide boundaries.

        Args:
            slide_index: Index of the slide to add table to (0-based)
            title: Title for the comparison table
            categories: List of comparison categories (row headers)
            option1_name: Name of first option (column header)
            option1_values: Values for first option
            option2_name: Name of second option (column header)
            option2_values: Values for second option
            option3_name: Optional name of third option
            option3_values: Optional values for third option
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            style: Table style (light, medium, dark)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming comparison table addition

        Example:
            await pptx_add_comparison_table(
                slide_index=2,
                title="Solution Comparison",
                categories=["Cost", "Time", "Features", "Support"],
                option1_name="Basic",
                option1_values=["$100/mo", "1 week", "Core features", "Email only"],
                option2_name="Professional",
                option2_values=["$500/mo", "2 days", "All features", "24/7 phone"],
                style="medium"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            # Build comparison table data
            headers = ["Category", option1_name, option2_name]
            if option3_name and option3_values:
                headers.append(option3_name)

            data = []
            for i, category in enumerate(categories):
                row = [category, option1_values[i], option2_values[i]]
                if option3_name and option3_values:
                    row.append(option3_values[i])
                data.append(row)

            slide = prs.slides[slide_index]

            # Get safe content area
            get_safe_content_area(has_title=bool(slide.shapes.title))

            # Validate position
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )

            # Remove overlapping placeholders (except title)
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, "shape_type"):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if hasattr(shape, "placeholder_format"):
                            from pptx.enum.shapes import PP_PLACEHOLDER

                            if shape.placeholder_format.type in [
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ]:
                                continue

                        shape_left = shape.left.inches if hasattr(shape.left, "inches") else 0
                        shape_top = shape.top.inches if hasattr(shape.top, "inches") else 0
                        shape_right = shape_left + (
                            shape.width.inches if hasattr(shape.width, "inches") else 0
                        )
                        shape_bottom = shape_top + (
                            shape.height.inches if hasattr(shape.height, "inches") else 0
                        )

                        if not (
                            shape_right < validated_left
                            or shape_left > validated_left + validated_width
                            or shape_bottom < validated_top
                            or shape_top > validated_top + validated_height
                        ):
                            placeholders_to_remove.append(shape)

            for placeholder in placeholders_to_remove:
                slide.shapes._spTree.remove(placeholder.element)

            # Get current theme
            theme_manager = ThemeManager()
            theme = theme_manager.get_theme(metadata.theme) if metadata.theme else None

            # Map style to variant
            variant_map = {"light": "minimal", "medium": "default", "dark": "bordered"}
            variant = variant_map.get(style, "default")

            # Create and render table using new Table component
            table_comp = Table(headers=headers, data=data, variant=variant, size="md", theme=theme)

            table_comp.render(
                slide,
                left=validated_left,
                top=validated_top,
                width=validated_width,
                height=validated_height,
            )

            # Update in VFS
            await manager.update(presentation)

            num_options = 3 if option3_name else 2
            message = f"Added {num_options}-way comparison table '{title}' to slide {slide_index}"
            return ComponentResponse(
                presentation=metadata.name,
                slide_index=slide_index,
                component="comparison_table",
                message=message,
                variant=None,
            ).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=f"Error adding comparison table: {str(e)}").model_dump_json()

    @mcp.tool
    async def pptx_update_table_cell(
        slide_index: int,
        table_index: int,
        row: int,
        col: int,
        new_value: str,
        bold: bool | None = None,
        color: str | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Update a specific cell in an existing table.

        Modifies the text and optionally the formatting of a table cell.

        Args:
            slide_index: Index of the slide containing the table (0-based)
            table_index: Index of the table on the slide (0-based)
            row: Row index of the cell (0-based)
            col: Column index of the cell (0-based)
            new_value: New text value for the cell
            bold: Whether to make text bold
            color: Hex color for text (e.g., "#FF0000")
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming cell update

        Example:
            await pptx_update_table_cell(
                slide_index=1,
                table_index=0,
                row=2,
                col=3,
                new_value="$150,000",
                bold=True,
                color="#008000"  # Green for positive numbers
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            slide = prs.slides[slide_index]

            # Find tables on the slide
            tables = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables.append(shape)

            if table_index >= len(tables):
                return ErrorResponse(
                    error=f"Table index {table_index} out of range. Found {len(tables)} tables on slide"
                ).model_dump_json()

            table = tables[table_index].table

            if row >= len(table.rows):
                return ErrorResponse(
                    error=f"Row {row} out of range. Table has {len(table.rows)} rows"
                ).model_dump_json()

            if col >= len(table.columns):
                return ErrorResponse(
                    error=f"Column {col} out of range. Table has {len(table.columns)} columns"
                ).model_dump_json()

            # Update the cell
            cell = table.cell(row, col)
            cell.text = new_value

            # Apply formatting if specified
            if bold is not None or color:
                paragraph = cell.text_frame.paragraphs[0]
                if bold is not None:
                    paragraph.font.bold = bold
                if color:
                    c_color = color[1:] if color.startswith("#") else color
                    try:
                        r = int(c_color[0:2], 16)
                        g = int(c_color[2:4], 16)
                        b = int(c_color[4:6], 16)
                        paragraph.font.color.rgb = RGBColor(r, g, b)
                    except (ValueError, IndexError):
                        pass

            # Update in VFS
            await manager.update(presentation)

            message = f"Updated table cell [{row}, {col}] to '{new_value}' on slide {slide_index}"
            return SuccessResponse(message=message).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=f"Error updating table cell: {str(e)}").model_dump_json()

    @mcp.tool
    async def pptx_format_table(
        slide_index: int,
        table_index: int,
        header_bold: bool = True,
        header_color: str | None = None,
        alternate_rows: bool = False,
        border_width: float | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Apply formatting to an entire table.

        Formats headers, alternating rows, borders, and other table styling.

        Args:
            slide_index: Index of the slide containing the table (0-based)
            table_index: Index of the table on the slide (0-based)
            header_bold: Whether to make header row bold
            header_color: Background color for header row (hex)
            alternate_rows: Whether to alternate row colors
            border_width: Width of table borders in points
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message with formatting applied

        Example:
            await pptx_format_table(
                slide_index=1,
                table_index=0,
                header_bold=True,
                header_color="#003366",
                alternate_rows=True,
                border_width=1.0
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            slide = prs.slides[slide_index]

            # Find tables
            tables = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables.append(shape)

            if table_index >= len(tables):
                return ErrorResponse(
                    error=f"Table index {table_index} out of range. Found {len(tables)} tables on slide"
                ).model_dump_json()

            table = tables[table_index].table
            formatting_applied = []

            # Format header row
            if len(table.rows) > 0:
                header_row = table.rows[0]
                for cell in header_row.cells:
                    if header_bold:
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True

                    if header_color:
                        h_color = header_color[1:] if header_color.startswith("#") else header_color
                        try:
                            r = int(h_color[0:2], 16)
                            g = int(h_color[2:4], 16)
                            b = int(h_color[4:6], 16)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(r, g, b)
                        except (ValueError, IndexError):
                            pass

                formatting_applied.append("header formatting")

            # Apply alternating row colors
            if alternate_rows:
                for i, row in enumerate(table.rows):
                    if i > 0 and i % 2 == 0:  # Skip header, alternate from row 2
                        for cell in row.cells:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
                formatting_applied.append("alternating rows")

            # Set border width
            if border_width:
                for row in table.rows:
                    for cell in row.cells:
                        cell.border_width = Pt(border_width)
                formatting_applied.append(f"border width {border_width}pt")

            # Update in VFS
            await manager.update(presentation)

            if formatting_applied:
                message = f"Applied table formatting: {', '.join(formatting_applied)}"
            else:
                message = "No formatting changes applied"

            return SuccessResponse(message=message).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=f"Error formatting table: {str(e)}").model_dump_json()

    # Return the tools for external access
    return {
        "pptx_add_data_table": pptx_add_data_table,
        "pptx_add_comparison_table": pptx_add_comparison_table,
        "pptx_update_table_cell": pptx_update_table_cell,
        "pptx_format_table": pptx_format_table,
    }
