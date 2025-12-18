"""
Image Tools for PowerPoint MCP Server

Provides async MCP tools for handling images in presentations.
Supports file paths, base64 data URLs, and various image formats.
"""

import base64
import io
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..layout.helpers import (
    validate_position,
    calculate_grid_layout,
    get_logo_position,
    get_safe_content_area,
    SLIDE_HEIGHT,
)
from ..constants import (
    SlideLayoutIndex,
    ErrorMessages,
)

# Import design system typography tokens
from ..tokens.typography import FONT_SIZES


def add_image(
    slide,
    image_source: str,
    left: float,
    top: float,
    width: float | None = None,
    height: float | None = None,
    maintain_ratio: bool = True,
):
    """
    Add an image to a slide.

    Helper function for image tools.

    Args:
        slide: Slide to add image to
        image_source: Path to image file or base64 data
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional, maintains ratio if not specified)
        height: Height in inches (optional, maintains ratio if not specified)
        maintain_ratio: Whether to maintain aspect ratio

    Returns:
        The created picture shape
    """
    # Handle base64 image data
    if image_source.startswith("data:image/"):
        header, encoded = image_source.split(",", 1)
        image_data = base64.b64decode(encoded)
        image_stream = io.BytesIO(image_data)

        if width and height:
            pic = slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top), width=Inches(width), height=Inches(height)
            )
        elif width:
            pic = slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top), width=Inches(width)
            )
        elif height:
            pic = slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top), height=Inches(height)
            )
        else:
            pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top))
    # Handle file path
    elif Path(image_source).exists():
        if width and height:
            pic = slide.shapes.add_picture(
                str(image_source),
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height),
            )
        elif width:
            pic = slide.shapes.add_picture(
                str(image_source), Inches(left), Inches(top), width=Inches(width)
            )
        elif height:
            pic = slide.shapes.add_picture(
                str(image_source), Inches(left), Inches(top), height=Inches(height)
            )
        else:
            pic = slide.shapes.add_picture(str(image_source), Inches(left), Inches(top))
    else:
        raise FileNotFoundError(f"Image file not found: {image_source}")

    return pic


def add_text_box_with_style(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    style_preset: str = "body",
):
    """
    Add a text box with a style preset.

    Helper function for adding styled text boxes.

    Args:
        slide: Slide to add text box to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        style_preset: Style preset ("caption", "body", "title", etc.)

    Returns:
        The created text box shape
    """
    from pptx.enum.text import PP_ALIGN

    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    text_frame = textbox.text_frame
    text_frame.text = text

    # Apply style based on preset using design tokens
    if style_preset == "caption":
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(FONT_SIZES["xs"])  # 10pt
            paragraph.font.italic = True
            paragraph.alignment = PP_ALIGN.CENTER
    elif style_preset == "title":
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(FONT_SIZES["3xl"])  # 28pt (was 24pt, adjusted to design token)
            paragraph.font.bold = True
    else:  # body or default
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(FONT_SIZES["sm"])  # 12pt

    return textbox


def register_image_tools(mcp, manager):
    """Register all image tools with the MCP server."""

    @mcp.tool
    async def pptx_add_image_slide(
        title: str, image_path: str, presentation: str | None = None
    ) -> str:
        """
        Add a slide with title and image.

        Creates a slide with a title and an image. The image can be provided
        as a file path or as base64-encoded data URL.

        Args:
            title: Title text for the slide
            image_path: Either a file path to an image or a data URL
                       (e.g., "data:image/png;base64,...")
            presentation: Name of presentation to add slide to (uses current if not specified)

        Returns:
            Success message confirming slide addition, or error message if image cannot be loaded

        Example:
            # Using file path
            await pptx_add_image_slide(
                title="Product Screenshot",
                image_path="/path/to/screenshot.png"
            )

            # Using base64 data URL
            await pptx_add_image_slide(
                title="Generated Chart",
                image_path="data:image/png;base64,iVBORw0KGgoAAAANS..."
            )
        """

        async def _add_image_slide():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            slide_layout = prs.slide_layouts[SlideLayoutIndex.TITLE_ONLY]  # Title only layout
            slide = prs.slides.add_slide(slide_layout)

            slide.shapes.title.text = title

            # Apply presentation theme to the slide first (even if image loading fails)
            if metadata and metadata.theme:
                from ..themes.theme_manager import ThemeManager

                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_theme(metadata.theme)
                if theme_obj:
                    theme_obj.apply_to_slide(slide)

            try:
                # For title-only layout, ensure image doesn't overlap with title
                # Title typically ends at 1.5 inches (0.3 start + 1.2 height)
                safe_top = 1.6  # Give a small margin below title
                left = Inches(0.5)
                top = Inches(safe_top)

                # Calculate appropriate width and height
                max_width = 9.0  # Leave margins on sides
                # Available height from below title to bottom with margin
                max_height = SLIDE_HEIGHT - safe_top - 0.5

                if image_path.startswith("data:image/"):
                    # Handle base64 image data URL
                    header, encoded = image_path.split(",", 1)
                    image_data = base64.b64decode(encoded)
                    image_stream = io.BytesIO(image_data)
                    # Use max_width but maintain aspect ratio
                    pic = slide.shapes.add_picture(image_stream, left, top, width=Inches(max_width))
                    # Check if height exceeds max and adjust if needed
                    if pic.height > Inches(max_height):
                        pic.height = Inches(max_height)
                        pic.width = Inches(max_height * max_width / max_height)  # Maintain ratio
                    result = f"Added image slide with base64 data to '{presentation or manager.get_current_name()}'"
                elif Path(image_path).exists():
                    # Handle file path
                    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(max_width))
                    # Check if height exceeds max and adjust if needed
                    if pic.height > Inches(max_height):
                        pic.height = Inches(max_height)
                        pic.width = Inches(max_height * max_width / max_height)  # Maintain ratio
                    result = f"Added image slide to '{presentation or manager.get_current_name()}'"
                else:
                    # Image not found - update VFS with themed slide anyway
                    await manager.update(presentation)
                    return f"Error: Image file not found: {image_path}"

                # Update in VFS if enabled
                await manager.update(presentation)
                return result

            except Exception as e:
                # Update VFS with themed slide even on error
                await manager.update(presentation)
                return f"Error: Failed to add image: {str(e)}"

        return await _add_image_slide()

    @mcp.tool
    async def pptx_add_image(
        slide_index: int,
        image_path: str,
        left: float = 1.0,
        top: float = 2.0,
        width: float | None = None,
        height: float | None = None,
        maintain_ratio: bool = True,
        presentation: str | None = None,
    ) -> str:
        """
        Add an image to an existing slide.

        Adds an image to a specific slide with precise positioning and sizing control.
        Supports both file paths and base64 data URLs.

        Args:
            slide_index: Index of the slide to add image to (0-based)
            image_path: Path to image file or base64 data URL
            left: Left position in inches
            top: Top position in inches
            width: Width in inches (optional, maintains ratio if not specified)
            height: Height in inches (optional, maintains ratio if not specified)
            maintain_ratio: Whether to maintain aspect ratio
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming image addition

        Example:
            await pptx_add_image(
                slide_index=1,
                image_path="/path/to/logo.png",
                left=0.5,
                top=0.5,
                width=2.0,
                height=1.0
            )
        """

        async def _add_image():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                # Validate and adjust position to ensure it fits within slide
                validated_left, validated_top, validated_width, validated_height = (
                    validate_position(
                        left,
                        top,
                        width or 3.0,  # Default width if not specified
                        height or 2.25,  # Default height if not specified
                    )
                )

                # Remove any overlapping placeholders (except title)
                placeholders_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, "shape_type"):
                        # Check if it's a placeholder (but not a title)
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            # Skip title placeholders
                            if hasattr(shape, "placeholder_format"):
                                from pptx.enum.shapes import PP_PLACEHOLDER

                                if shape.placeholder_format.type in [
                                    PP_PLACEHOLDER.TITLE,
                                    PP_PLACEHOLDER.CENTER_TITLE,
                                ]:
                                    continue

                            # Check if placeholder overlaps with image area
                            shape_left = shape.left.inches if hasattr(shape.left, "inches") else 0
                            shape_top = shape.top.inches if hasattr(shape.top, "inches") else 0
                            shape_right = shape_left + (
                                shape.width.inches if hasattr(shape.width, "inches") else 0
                            )
                            shape_bottom = shape_top + (
                                shape.height.inches if hasattr(shape.height, "inches") else 0
                            )

                            img_right = validated_left + validated_width
                            img_bottom = validated_top + validated_height

                            # Check for overlap
                            if not (
                                shape_right < validated_left
                                or shape_left > img_right
                                or shape_bottom < validated_top
                                or shape_top > img_bottom
                            ):
                                placeholders_to_remove.append(shape)

                # Remove overlapping placeholders
                for placeholder in placeholders_to_remove:
                    slide.shapes._spTree.remove(placeholder.element)

                # Use validated dimensions
                add_image(
                    slide,
                    image_path,
                    validated_left,
                    validated_top,
                    validated_width if width else None,
                    validated_height if height else None,
                    maintain_ratio,
                )

                # Update in VFS if enabled
                await manager.update(presentation)

                # Report if position was adjusted
                position_note = ""
                if (
                    width
                    and height
                    and (
                        validated_left != left
                        or validated_top != top
                        or validated_width != width
                        or validated_height != height
                    )
                ):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"

                return f"Added image to slide {slide_index}{position_note}"
            except Exception as e:
                return f"Error adding image: {str(e)}"

        return await _add_image()

    @mcp.tool
    async def pptx_add_background_image(
        slide_index: int, image_path: str, presentation: str | None = None
    ) -> str:
        """
        Set a background image for a slide.

        Sets an image as the background of a specific slide, covering the entire slide area.

        Args:
            slide_index: Index of the slide to set background for (0-based)
            image_path: Path to image file or base64 data URL
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming background addition

        Example:
            await pptx_add_background_image(
                slide_index=0,
                image_path="/path/to/background.jpg"
            )
        """

        async def _add_background():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                # Add image covering entire slide (10x7.5 inches is standard slide size)
                pic_shape = add_image(slide, image_path, 0, 0, 10, 7.5, maintain_ratio=False)

                # Send image to back so it acts as background
                pic_shape.element.getparent().remove(pic_shape.element)
                slide.shapes._spTree.insert(
                    2, pic_shape.element
                )  # Insert after slide layout elements

                # Update in VFS if enabled
                await manager.update(presentation)

                return f"Set background image for slide {slide_index}"
            except Exception as e:
                return f"Error setting background image: {str(e)}"

        return await _add_background()

    @mcp.tool
    async def pptx_add_image_gallery(
        slide_index: int,
        image_paths: list[str],
        columns: int = 2,
        spacing: float = 0.2,
        start_left: float = 1.0,
        start_top: float = 2.0,
        image_width: float = 3.0,
        image_height: float = 2.25,
        presentation: str | None = None,
    ) -> str:
        """
        Add multiple images in a grid layout to a slide.

        Creates a photo gallery or image grid on a slide with automatic positioning.

        Args:
            slide_index: Index of the slide to add images to (0-based)
            image_paths: List of image paths or base64 data URLs
            columns: Number of columns in the grid
            spacing: Spacing between images in inches
            start_left: Starting left position in inches
            start_top: Starting top position in inches
            image_width: Width of each image in inches
            image_height: Height of each image in inches
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message with count of images added

        Example:
            await pptx_add_image_gallery(
                slide_index=2,
                image_paths=[
                    "/path/to/photo1.jpg",
                    "/path/to/photo2.jpg",
                    "/path/to/photo3.jpg",
                    "/path/to/photo4.jpg"
                ],
                columns=2,
                spacing=0.3
            )
        """

        async def _add_gallery():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                added_count = 0

                # Remove content placeholders that might interfere with gallery
                placeholders_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, "shape_type"):
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            # Skip title placeholders
                            if hasattr(shape, "placeholder_format"):
                                from pptx.enum.shapes import PP_PLACEHOLDER

                                if shape.placeholder_format.type in [
                                    PP_PLACEHOLDER.TITLE,
                                    PP_PLACEHOLDER.CENTER_TITLE,
                                ]:
                                    continue
                            # Remove content placeholders
                            placeholders_to_remove.append(shape)

                for placeholder in placeholders_to_remove:
                    slide.shapes._spTree.remove(placeholder.element)

                # Calculate grid layout with validation
                safe_area = get_safe_content_area(has_title=True)
                positions = calculate_grid_layout(
                    num_items=len(image_paths),
                    columns=columns,
                    spacing=spacing,
                    container_left=start_left,
                    container_top=start_top,
                    container_width=safe_area["width"],
                    container_height=safe_area["height"] - (start_top - safe_area["top"]),
                )

                for i, (image_path, pos) in enumerate(zip(image_paths, positions)):
                    try:
                        add_image(
                            slide,
                            image_path,
                            pos["left"],
                            pos["top"],
                            pos["width"],
                            pos["height"],
                            maintain_ratio=False,
                        )
                        added_count += 1
                    except Exception as img_error:
                        # Continue with other images if one fails
                        print(f"Warning: Failed to add image {image_path}: {img_error}")

                # Update in VFS if enabled
                await manager.update(presentation)

                return f"Added {added_count} of {len(image_paths)} images to gallery on slide {slide_index}"
            except Exception as e:
                return f"Error creating image gallery: {str(e)}"

        return await _add_gallery()

    @mcp.tool
    async def pptx_add_image_with_caption(
        slide_index: int,
        image_path: str,
        caption: str,
        left: float = 1.0,
        top: float = 1.5,
        image_width: float = 6.0,
        image_height: float = 4.0,
        caption_height: float = 0.8,
        presentation: str | None = None,
    ) -> str:
        """
        Add an image with a caption text box below it.

        Creates an image with an accompanying caption text box positioned below the image.

        Args:
            slide_index: Index of the slide to add image to (0-based)
            image_path: Path to image file or base64 data URL
            caption: Caption text to display below the image
            left: Left position in inches for both image and caption
            top: Top position in inches for the image
            image_width: Width of the image in inches
            image_height: Height of the image in inches
            caption_height: Height of the caption text box in inches
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming image and caption addition

        Example:
            await pptx_add_image_with_caption(
                slide_index=3,
                image_path="/path/to/product.jpg",
                caption="Our flagship product - Model X Pro",
                left=2.0,
                top=1.0,
                image_width=6.0,
                image_height=4.0
            )
        """

        async def _add_image_caption():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                # Validate position for image and caption
                validated_left, validated_top, validated_width, validated_height = (
                    validate_position(left, top, image_width, image_height + caption_height + 0.1)
                )

                # Adjust image height if needed to fit caption
                img_height = validated_height - caption_height - 0.1

                # Remove any overlapping placeholders (except title)
                placeholders_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, "shape_type"):
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            # Skip title placeholders
                            if hasattr(shape, "placeholder_format"):
                                from pptx.enum.shapes import PP_PLACEHOLDER

                                if shape.placeholder_format.type in [
                                    PP_PLACEHOLDER.TITLE,
                                    PP_PLACEHOLDER.CENTER_TITLE,
                                ]:
                                    continue

                            shape_left = shape.left.inches if hasattr(shape.left, "inches") else 0
                            shape_top = shape.top.inches if hasattr(shape.top, "inches") else 0
                            shape_right = shape_left + (
                                shape.width.inches if hasattr(shape.width, "inches") else 0
                            )
                            shape_bottom = shape_top + (
                                shape.height.inches if hasattr(shape.height, "inches") else 0
                            )

                            # Check overlap with combined image+caption area
                            if not (
                                shape_right < validated_left
                                or shape_left > validated_left + validated_width
                                or shape_bottom < validated_top
                                or shape_top > validated_top + validated_height
                            ):
                                placeholders_to_remove.append(shape)

                for placeholder in placeholders_to_remove:
                    slide.shapes._spTree.remove(placeholder.element)

                # Add the image
                add_image(
                    slide,
                    image_path,
                    validated_left,
                    validated_top,
                    validated_width,
                    img_height,
                    maintain_ratio=False,
                )

                # Add caption text box below the image
                caption_top = (
                    validated_top + img_height + 0.1
                )  # Small gap between image and caption

                add_text_box_with_style(
                    slide,
                    validated_left,
                    caption_top,
                    validated_width,
                    caption_height,
                    caption,
                    style_preset="caption",
                )

                # Update in VFS if enabled
                await manager.update(presentation)

                # Report if position was adjusted
                position_note = ""
                total_height = image_height + caption_height + 0.1
                if (
                    validated_left != left
                    or validated_top != top
                    or validated_width != image_width
                    or validated_height != total_height
                ):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"

                return f"Added image with caption to slide {slide_index}{position_note}"
            except Exception as e:
                return f"Error adding image with caption: {str(e)}"

        return await _add_image_caption()

    @mcp.tool
    async def pptx_add_logo(
        slide_index: int,
        logo_path: str,
        position: str = "top-right",
        size: float = 1.0,
        margin: float = 0.3,
        presentation: str | None = None,
    ) -> str:
        """
        Add a logo to a slide in a standard position.

        Adds a company logo or brand image to a slide in common positions
        like corners or along edges.

        Args:
            slide_index: Index of the slide to add logo to (0-based)
            logo_path: Path to logo image file or base64 data URL
            position: Logo position - "top-left", "top-right", "bottom-left",
                     "bottom-right", "center", "top-center", "bottom-center"
            size: Size of the logo in inches (width, height maintains ratio)
            margin: Margin from slide edges in inches
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming logo addition

        Example:
            await pptx_add_logo(
                slide_index=0,
                logo_path="/path/to/company-logo.png",
                position="top-right",
                size=1.5,
                margin=0.4
            )
        """

        async def _add_logo():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                # Use layout helper for logo positioning
                # Note: get_logo_position defaults to top-right for invalid positions
                logo_pos = get_logo_position(position, size, margin)

                left = logo_pos["left"]
                top = logo_pos["top"]

                # Add logo with maintained aspect ratio
                add_image(slide, logo_path, left, top, size, None, maintain_ratio=True)

                # Update in VFS if enabled
                await manager.update(presentation)

                return f"Added logo to {position} of slide {slide_index}"
            except Exception as e:
                return f"Error adding logo: {str(e)}"

        return await _add_logo()

    @mcp.tool
    async def pptx_replace_image(
        slide_index: int,
        old_image_index: int,
        new_image_path: str,
        maintain_position: bool = True,
        maintain_size: bool = True,
        presentation: str | None = None,
    ) -> str:
        """
        Replace an existing image on a slide with a new image.

        Replaces an image while optionally maintaining its position and size.

        Args:
            slide_index: Index of the slide containing the image (0-based)
            old_image_index: Index of the image to replace (0-based, order they appear on slide)
            new_image_path: Path to new image file or base64 data URL
            maintain_position: Whether to keep the same position
            maintain_size: Whether to keep the same size
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming image replacement

        Example:
            await pptx_replace_image(
                slide_index=1,
                old_image_index=0,
                new_image_path="/path/to/new_chart.png",
                maintain_position=True,
                maintain_size=True
            )
        """

        async def _replace_image():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                # Find image shapes on the slide
                image_shapes = []
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image_shapes.append(shape)

                if old_image_index >= len(image_shapes):
                    return f"Error: Image index {old_image_index} out of range. Found {len(image_shapes)} images on slide."

                old_image = image_shapes[old_image_index]

                # Get current position and size
                if maintain_position:
                    left = old_image.left.inches
                    top = old_image.top.inches
                else:
                    # Use safe content area for default position
                    safe_area = get_safe_content_area(has_title=True)
                    left, top = safe_area["left"], safe_area["top"]

                if maintain_size:
                    width = old_image.width.inches
                    height = old_image.height.inches
                else:
                    width, height = None, None  # Let it size automatically

                # Store original values for reporting
                orig_left, orig_top, orig_width, orig_height = left, top, width, height

                # Validate position if specified
                if width and height:
                    validated_left, validated_top, validated_width, validated_height = (
                        validate_position(left, top, width, height)
                    )
                else:
                    validated_left, validated_top, validated_width, validated_height = (
                        left,
                        top,
                        width,
                        height,
                    )

                # Remove old image
                slide.shapes._spTree.remove(old_image.element)

                # Add new image
                add_image(
                    slide,
                    new_image_path,
                    validated_left,
                    validated_top,
                    validated_width,
                    validated_height,
                    maintain_ratio=not maintain_size,
                )

                # Update in VFS if enabled
                await manager.update(presentation)

                # Report if position was adjusted
                position_note = ""
                if (
                    width
                    and height
                    and (
                        validated_left != orig_left
                        or validated_top != orig_top
                        or validated_width != orig_width
                        or validated_height != orig_height
                    )
                ):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"

                return f"Replaced image {old_image_index} on slide {slide_index}{position_note}"
            except Exception as e:
                return f"Error replacing image: {str(e)}"

        return await _replace_image()

    @mcp.tool
    async def pptx_add_image_placeholder(
        slide_index: int,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 4.0,
        height: float = 3.0,
        label: str = "Image Placeholder",
        background_color: str = "#E0E0E0",
        text_color: str = "#666666",
        presentation: str | None = None,
    ) -> str:
        """
        Add an image placeholder for mockups and prototyping.

        Creates a rectangular placeholder box with a label, useful for designing
        slide layouts before actual images are available.

        Args:
            slide_index: Index of the slide to add placeholder to (0-based)
            left: Left position in inches
            top: Top position in inches
            width: Width of placeholder in inches
            height: Height of placeholder in inches
            label: Text to display on the placeholder
            background_color: Background color as hex string (default: light gray)
            text_color: Text color as hex string (default: dark gray)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming placeholder addition

        Example:
            await pptx_add_image_placeholder(
                slide_index=1,
                left=2.0,
                top=2.0,
                width=5.0,
                height=3.5,
                label="Product Screenshot",
                background_color="#F5F5F5"
            )
        """

        async def _add_placeholder():
            result = await manager.get(presentation)
            if not result:
                return ErrorMessages.NO_PRESENTATION

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            try:
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.dml.color import RGBColor
                from pptx.util import Pt
                from pptx.enum.text import PP_ALIGN

                # Validate position
                validated_left, validated_top, validated_width, validated_height = (
                    validate_position(left, top, width, height)
                )

                # Create rectangle shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(validated_left),
                    Inches(validated_top),
                    Inches(validated_width),
                    Inches(validated_height),
                )

                # Parse and apply background color
                if background_color.startswith("#"):
                    bg_hex = background_color[1:]
                else:
                    bg_hex = background_color

                try:
                    r, g, b = int(bg_hex[0:2], 16), int(bg_hex[2:4], 16), int(bg_hex[4:6], 16)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    # Fallback to light gray if color parsing fails
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(224, 224, 224)

                # Add border
                shape.line.color.rgb = RGBColor(150, 150, 150)
                shape.line.width = Pt(1)
                shape.line.dash_style = 2  # Dashed line

                # Add label text
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = label
                p.alignment = PP_ALIGN.CENTER

                # Parse and apply text color
                if text_color.startswith("#"):
                    txt_hex = text_color[1:]
                else:
                    txt_hex = text_color

                try:
                    r, g, b = int(txt_hex[0:2], 16), int(txt_hex[2:4], 16), int(txt_hex[4:6], 16)
                    p.font.color.rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    # Fallback to dark gray if color parsing fails
                    p.font.color.rgb = RGBColor(102, 102, 102)

                p.font.size = Pt(FONT_SIZES["base"])  # 14pt
                p.font.italic = True

                # Center text vertically
                text_frame.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE

                # Update in VFS if enabled
                await manager.update(presentation)

                # Report if position was adjusted
                position_note = ""
                if (
                    validated_left != left
                    or validated_top != top
                    or validated_width != width
                    or validated_height != height
                ):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"

                return f"Added image placeholder '{label}' to slide {slide_index}{position_note}"
            except Exception as e:
                return f"Error adding image placeholder: {str(e)}"

        return await _add_placeholder()

    # Return the tools for external access
    return {
        "pptx_add_image_slide": pptx_add_image_slide,
        "pptx_add_image": pptx_add_image,
        "pptx_add_background_image": pptx_add_background_image,
        "pptx_add_image_gallery": pptx_add_image_gallery,
        "pptx_add_image_with_caption": pptx_add_image_with_caption,
        "pptx_add_logo": pptx_add_logo,
        "pptx_replace_image": pptx_replace_image,
        "pptx_add_image_placeholder": pptx_add_image_placeholder,
    }
