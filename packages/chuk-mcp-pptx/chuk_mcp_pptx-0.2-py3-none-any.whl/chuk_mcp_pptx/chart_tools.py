"""
Chart Tools for PowerPoint MCP Server

Provides a unified async MCP tool for creating all chart types in presentations.
Optimized for AI/LLM tool invocation with clear parameter structure.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE


def register_chart_tools(mcp, manager):
    """Register the unified chart tool with the MCP server."""

    from .utilities.chart_utils import add_chart, add_pie_chart, add_scatter_chart
    from .layout.helpers import validate_position, get_safe_content_area
    from .tokens.colors import PALETTE
    import json

    @mcp.tool
    async def pptx_add_chart(
        slide_index: int,
        chart_type: str,
        data: dict,
        title: str = "",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.5,
        options: dict | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Create and add a chart (bar, line, pie, etc.) to a presentation slide.

        USE THIS TOOL when you need to add visual data charts to slides.
        This is the PRIMARY tool for creating charts - use it for all chart types.

        Common use cases:
        - "add a bar chart showing sales data" → use this tool
        - "create a pie chart for market share" → use this tool
        - "show revenue trends in a line chart" → use this tool

        ⚠️  BEST PRACTICES:
        - Multiple charts: Use SEPARATE SLIDES for clarity (1 chart/slide ideal, 2 max)
        - Chart titles: Keep concise (under 45 characters) to prevent wrapping/overlap
        - THREE+ charts = create separate slides to avoid overlap/clutter

        Args:
            slide_index: Index of the slide to add chart to (0-based)

            chart_type: Type of chart to create. Options:
                - "column" - Vertical bars comparing values
                - "column_stacked" - Stacked vertical bars
                - "bar" - Horizontal bars
                - "bar_stacked" - Stacked horizontal bars
                - "line" - Line graph showing trends
                - "line_markers" - Line graph with data point markers
                - "pie" - Pie chart showing proportions
                - "doughnut" - Doughnut chart (pie with hollow center)
                - "area" - Area chart showing magnitude over time
                - "area_stacked" - Stacked area chart
                - "scatter" - XY scatter plot for correlations
                - "bubble" - Bubble chart for 3D data (x, y, size)
                - "radar" - Radar/spider chart for multi-criteria
                - "radar_filled" - Filled radar chart
                - "waterfall" - Waterfall chart for incremental changes

            data: Chart data structure. Format depends on chart_type:

                For column/bar/line/area charts:
                {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": {
                        "Revenue": [100, 120, 140, 160],
                        "Profit": [20, 25, 30, 35]
                    }
                }

                For pie/doughnut charts:
                {
                    "categories": ["Product A", "Product B", "Product C"],
                    "values": [45, 30, 25]
                }

                For scatter charts:
                {
                    "series": [
                        {
                            "name": "Dataset 1",
                            "x_values": [1, 2, 3, 4, 5],
                            "y_values": [2, 4, 6, 8, 10]
                        }
                    ]
                }

                For bubble charts:
                {
                    "series": [
                        {
                            "name": "Markets",
                            "points": [[10, 20, 5], [15, 25, 8], [20, 30, 12]]
                        }
                    ]
                }

                For radar charts:
                {
                    "categories": ["Speed", "Reliability", "Comfort", "Design"],
                    "series": {
                        "Model A": [8, 7, 9, 8],
                        "Model B": [7, 9, 7, 6]
                    }
                }

                For waterfall charts:
                {
                    "categories": ["Start", "Sales", "Costs", "Tax", "End"],
                    "values": [100, 50, -30, -10, 110]
                }

            title: Optional chart title
            left: Left position in inches (will be validated)
            top: Top position in inches (will be validated)
            width: Width in inches (will be validated)
            height: Height in inches (will be validated)

            options: Optional chart-specific options:
                {
                    "show_percentages": true,  # For pie charts
                    "show_legend": true,        # Show/hide legend
                    "legend_position": "right", # Legend position: right, left, top, bottom
                    "colors": ["#FF5733", "#33FF57", "#3357FF"]  # Custom colors
                }

            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming chart addition, or error with guidance

        Examples:
            # Column chart
            await pptx_add_chart(
                slide_index=1,
                chart_type="column",
                data={
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": {
                        "Revenue": [100, 120, 140, 160],
                        "Profit": [20, 25, 30, 35]
                    }
                },
                title="Quarterly Performance"
            )

            # Pie chart
            await pptx_add_chart(
                slide_index=2,
                chart_type="pie",
                data={
                    "categories": ["Product A", "Product B", "Product C"],
                    "values": [45, 30, 25]
                },
                title="Market Share",
                options={"show_percentages": True}
            )

            # Scatter plot
            await pptx_add_chart(
                slide_index=3,
                chart_type="scatter",
                data={
                    "series": [
                        {
                            "name": "Sales Data",
                            "x_values": [10, 20, 30, 40, 50],
                            "y_values": [15, 25, 45, 35, 55]
                        }
                    ]
                },
                title="Price vs Sales Correlation"
            )
        """

        async def _add_unified_chart():
            nonlocal options, data

            # Convert parameters to correct types (MCP protocol may send as strings)
            try:
                idx = int(slide_index)
            except (ValueError, TypeError):
                return f"Error: slide_index must be a number, got: {slide_index}"

            try:
                # Store original values for comparison
                original_left = float(left)
                original_top = float(top)
                original_width = float(width)
                original_height = float(height)

                # Use these for validation
                validated_left = original_left
                validated_top = original_top
                validated_width = original_width
                validated_height = original_height
            except (ValueError, TypeError):
                return f"Error: Position/size parameters must be numbers. Got left={left}, top={top}, width={width}, height={height}"

            # Parse data if it's a JSON string
            if isinstance(data, str):
                try:
                    data = json.loads(data)
                except json.JSONDecodeError as e:
                    return (
                        f"Error: data parameter must be a valid JSON object. Got parse error: {e}"
                    )

            # Parse options if it's a JSON string
            if isinstance(options, str):
                try:
                    options = json.loads(options)
                except json.JSONDecodeError:
                    options = {}

            result = await manager.get(presentation)
            if not result:
                return "Error: No presentation found. Create one first with pptx_create()"

            prs, metadata = result

            if idx >= len(prs.slides) or idx < 0:
                return f"Error: Slide index {idx} out of range. Presentation has {len(prs.slides)} slides (0-{len(prs.slides) - 1})."

            slide = prs.slides[idx]

            # Get safe content area
            safe_area = get_safe_content_area(has_title=bool(slide.shapes.title))

            # Validate and adjust position to fit within slide (use already converted float values)
            validated_left, validated_top, validated_width, validated_height = validate_position(
                validated_left, validated_top, validated_width, validated_height
            )

            # Further adjust if position is too close to title area
            if slide.shapes.title and validated_top < safe_area["top"]:
                validated_top = safe_area["top"]

            # Default options
            if options is None:
                options = {}

            # Remove any overlapping placeholders (except title)
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, "shape_type"):
                    # Check if it's a placeholder (but not a title)
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, "placeholder_format"):
                            from pptx.enum.shapes import PP_PLACEHOLDER

                            if shape.placeholder_format.type in [
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            ]:
                                continue

                        # Check if placeholder overlaps with chart area
                        shape_left = shape.left.inches if hasattr(shape.left, "inches") else 0
                        shape_top = shape.top.inches if hasattr(shape.top, "inches") else 0
                        shape_right = shape_left + (
                            shape.width.inches if hasattr(shape.width, "inches") else 0
                        )
                        shape_bottom = shape_top + (
                            shape.height.inches if hasattr(shape.height, "inches") else 0
                        )

                        chart_right = validated_left + validated_width
                        chart_bottom = validated_top + validated_height

                        # Check for overlap
                        if not (
                            shape_right < validated_left
                            or shape_left > chart_right
                            or shape_bottom < validated_top
                            or shape_top > chart_bottom
                        ):
                            placeholders_to_remove.append(shape)

            # Remove overlapping placeholders
            for placeholder in placeholders_to_remove:
                slide.shapes._spTree.remove(placeholder.element)

            # Check for potential chart overlap with existing charts
            overlap_warning = ""
            chart_right = validated_left + validated_width
            chart_bottom = validated_top + validated_height

            for shape in slide.shapes:
                if hasattr(shape, "has_chart") and shape.has_chart:
                    existing_left = shape.left.inches if hasattr(shape.left, "inches") else 0
                    existing_top = shape.top.inches if hasattr(shape.top, "inches") else 0
                    existing_right = existing_left + (
                        shape.width.inches if hasattr(shape.width, "inches") else 0
                    )
                    existing_bottom = existing_top + (
                        shape.height.inches if hasattr(shape.height, "inches") else 0
                    )

                    # Check for overlap
                    if not (
                        existing_right <= validated_left
                        or existing_left >= chart_right
                        or existing_bottom <= validated_top
                        or existing_top >= chart_bottom
                    ):
                        overlap_warning = " ⚠️  Warning: Chart may overlap with existing chart on this slide. Consider using separate slides for clarity."
                        break

            try:
                # Validate chart type
                valid_types = [
                    "column",
                    "column_stacked",
                    "bar",
                    "bar_stacked",
                    "line",
                    "line_markers",
                    "pie",
                    "doughnut",
                    "area",
                    "area_stacked",
                    "scatter",
                    "bubble",
                    "radar",
                    "radar_filled",
                    "waterfall",
                ]

                if chart_type not in valid_types:
                    return f"Error: Invalid chart_type '{chart_type}'. Valid types: {', '.join(valid_types)}"

                # Handle different chart types based on their data requirements
                if chart_type in [
                    "column",
                    "column_stacked",
                    "bar",
                    "bar_stacked",
                    "line",
                    "line_markers",
                    "area",
                    "area_stacked",
                ]:
                    # Standard category charts
                    if "categories" not in data or "series" not in data:
                        return f"Error: {chart_type} charts require 'categories' and 'series' in data. Got: {list(data.keys())}"

                    categories = data["categories"]
                    series_data = data["series"]

                    # Map radar_filled to the enum name
                    chart_type_mapped = chart_type

                    chart_shape = add_chart(
                        slide,
                        chart_type_mapped,
                        validated_left,
                        validated_top,
                        validated_width,
                        validated_height,
                        categories,
                        series_data,
                        title,
                        has_legend=options.get("show_legend", True),
                        legend_position=options.get("legend_position", "right"),
                    )

                elif chart_type in ["pie", "doughnut"]:
                    # Pie-type charts
                    if "categories" not in data or "values" not in data:
                        return f"Error: {chart_type} charts require 'categories' and 'values' in data. Got: {list(data.keys())}"

                    categories = data["categories"]
                    values = data["values"]

                    if chart_type == "pie":
                        chart_shape = add_pie_chart(
                            slide,
                            validated_left,
                            validated_top,
                            validated_width,
                            validated_height,
                            categories,
                            values,
                            title,
                            show_percentages=options.get("show_percentages", True),
                        )
                    else:  # doughnut
                        # Use the general add_chart for doughnut
                        series_data = {"Values": values}
                        chart_shape = add_chart(
                            slide,
                            "doughnut",
                            validated_left,
                            validated_top,
                            validated_width,
                            validated_height,
                            categories,
                            series_data,
                            title,
                        )

                elif chart_type == "scatter":
                    # Scatter charts
                    if "series" not in data:
                        return f"Error: Scatter charts require 'series' in data with 'x_values' and 'y_values'. Got: {list(data.keys())}"

                    series_data = data["series"]
                    chart_shape = add_scatter_chart(
                        slide,
                        validated_left,
                        validated_top,
                        validated_width,
                        validated_height,
                        series_data,
                        title,
                        has_legend=options.get("show_legend", True),
                    )

                elif chart_type == "bubble":
                    # Bubble charts
                    if "series" not in data:
                        return f"Error: Bubble charts require 'series' in data with 'points' as [x, y, size]. Got: {list(data.keys())}"

                    from pptx.chart.data import BubbleChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    from pptx.util import Inches, Pt

                    chart_data = BubbleChartData()

                    for series in data["series"]:
                        series_obj = chart_data.add_series(series.get("name", "Series"))
                        for point in series.get("points", []):
                            if len(point) == 3:
                                series_obj.add_data_point(point[0], point[1], point[2])
                            else:
                                return (
                                    f"Error: Bubble chart points must be [x, y, size]. Got: {point}"
                                )

                    chart_shape = slide.shapes.add_chart(
                        XL_CHART_TYPE.BUBBLE,
                        Inches(validated_left),
                        Inches(validated_top),
                        Inches(validated_width),
                        Inches(validated_height),
                        chart_data,
                    )

                    if title:
                        chart = chart_shape.chart
                        chart.has_title = True
                        chart.chart_title.text_frame.text = title
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
                        chart.chart_title.text_frame.paragraphs[0].font.bold = True

                elif chart_type in ["radar", "radar_filled"]:
                    # Radar charts
                    if "categories" not in data or "series" not in data:
                        return f"Error: Radar charts require 'categories' and 'series' in data. Got: {list(data.keys())}"

                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
                    from pptx.util import Inches, Pt

                    chart_data = CategoryChartData()
                    chart_data.categories = data["categories"]

                    for series_name, values in data["series"].items():
                        chart_data.add_series(series_name, values)

                    chart_enum = (
                        XL_CHART_TYPE.RADAR_FILLED
                        if chart_type == "radar_filled"
                        else XL_CHART_TYPE.RADAR
                    )

                    chart_shape = slide.shapes.add_chart(
                        chart_enum,
                        Inches(validated_left),
                        Inches(validated_top),
                        Inches(validated_width),
                        Inches(validated_height),
                        chart_data,
                    )

                    chart = chart_shape.chart
                    if title:
                        chart.has_title = True
                        chart.chart_title.text_frame.text = title
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
                        chart.chart_title.text_frame.paragraphs[0].font.bold = True

                    chart.has_legend = options.get("show_legend", True)
                    if chart.has_legend:
                        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                        chart.legend.include_in_layout = False

                elif chart_type == "waterfall":
                    # Waterfall charts
                    if "categories" not in data or "values" not in data:
                        return f"Error: Waterfall charts require 'categories' and 'values' in data. Got: {list(data.keys())}"

                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor

                    # Calculate cumulative values for waterfall effect
                    values = data["values"]
                    cumulative = []
                    running_total = 0
                    for val in values:
                        running_total += val
                        cumulative.append(running_total)

                    chart_data = CategoryChartData()
                    chart_data.categories = data["categories"]
                    chart_data.add_series("Values", cumulative)

                    chart_shape = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(validated_left),
                        Inches(validated_top),
                        Inches(validated_width),
                        Inches(validated_height),
                        chart_data,
                    )

                    chart = chart_shape.chart

                    if title:
                        chart.has_title = True
                        chart.chart_title.text_frame.text = title
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
                        chart.chart_title.text_frame.paragraphs[0].font.bold = True

                    # Color bars based on positive/negative values
                    series = chart.series[0]
                    for idx, val in enumerate(values):
                        point = series.points[idx]
                        fill = point.format.fill
                        fill.solid()
                        if val > 0:
                            fill.fore_color.rgb = RGBColor(0, 176, 80)  # Green
                        elif val < 0:
                            fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
                        else:
                            fill.fore_color.rgb = RGBColor(128, 128, 128)  # Gray

                    chart.has_legend = False  # No legend for waterfall

                else:
                    return f"Error: Unsupported chart_type '{chart_type}'"

                # Apply theme-aware colors to chart text elements
                if metadata and metadata.theme:
                    from .themes.theme_manager import ThemeManager

                    theme_manager = ThemeManager()
                    theme_obj = theme_manager.get_theme(metadata.theme)
                    if theme_obj:
                        # Only apply background, don't override chart text colors
                        theme_obj.apply_to_slide(slide, override_text_colors=False)

                        # Apply theme foreground color to chart text elements
                        text_color = theme_obj.get_color("foreground.DEFAULT")

                        # Get chart object from chart_shape
                        if hasattr(chart_shape, "chart"):
                            chart = chart_shape.chart

                            # Style chart title
                            if chart.has_title:
                                try:
                                    for para in chart.chart_title.text_frame.paragraphs:
                                        para.font.color.rgb = text_color
                                except (AttributeError, ValueError):
                                    pass

                            # Style axis labels
                            try:
                                if hasattr(chart, "value_axis") and hasattr(
                                    chart.value_axis, "tick_labels"
                                ):
                                    chart.value_axis.tick_labels.font.color.rgb = text_color
                            except (AttributeError, ValueError):
                                pass

                            try:
                                if hasattr(chart, "category_axis") and hasattr(
                                    chart.category_axis, "tick_labels"
                                ):
                                    chart.category_axis.tick_labels.font.color.rgb = text_color
                            except (AttributeError, ValueError):
                                pass

                            # Style legend
                            try:
                                if chart.has_legend and hasattr(chart.legend, "font"):
                                    chart.legend.font.color.rgb = text_color
                            except (AttributeError, ValueError):
                                pass

                            # Style data labels if present
                            try:
                                for series in chart.series:
                                    if (
                                        hasattr(series, "has_data_labels")
                                        and series.has_data_labels
                                    ):
                                        for point in series.points:
                                            if hasattr(point, "data_label"):
                                                point.data_label.font.color.rgb = text_color
                            except (AttributeError, ValueError):
                                pass

                # Update in VFS if enabled
                await manager.update(presentation)

                # Report if position was adjusted
                position_note = ""
                if (
                    validated_left != original_left
                    or validated_top != original_top
                    or validated_width != original_width
                    or validated_height != original_height
                ):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"

                return f"Added {chart_type} chart to slide {idx}{position_note}{overlap_warning}"

            except KeyError as e:
                return f"Error: Missing required data field: {str(e)}. Check the data structure for {chart_type} charts."
            except ValueError as e:
                return f"Error: Invalid data values: {str(e)}"
            except Exception as e:
                return f"Error adding {chart_type} chart: {str(e)}"

        return await _add_unified_chart()

    @mcp.tool
    async def pptx_get_chart_style(style_preset: str = "corporate") -> str:
        """
        Get a pre-defined color palette for charts based on a style preset.

        Returns a color palette from the design token system that can be used
        with the pptx_add_chart tool's options.colors parameter.

        Available style presets:
        - "corporate": Professional blues and grays for business presentations
        - "vibrant": Bold, energetic colors for creative presentations
        - "pastel": Soft, muted colors for gentle emphasis
        - "monochrome_blue": Shades of blue for focused data
        - "monochrome_green": Shades of green for growth/positive metrics
        - "earthy": Natural, warm tones (browns, greens, oranges)
        - "cool": Blues, cyans, and teals for calm data
        - "warm": Reds, oranges, and yellows for energy
        - "rainbow": Full spectrum for diverse categories
        - "status": Traffic light colors (green/yellow/red) for KPIs

        Args:
            style_preset: Name of the style preset to use

        Returns:
            JSON with color array that can be used with pptx_add_chart

        Example:
            style = await pptx_get_chart_style("corporate")
            # Returns: {"colors": ["#3b82f6", "#64748b", "#10b981", ...]}

            # Use with pptx_add_chart:
            await pptx_add_chart(
                slide_index=0,
                chart_type="column",
                data={"categories": [...], "series": {...}},
                options=json.loads(style)
            )
        """
        # Define style presets using colors from the token system
        style_presets = {
            "corporate": {
                "colors": [
                    PALETTE["blue"][600],  # Professional blue
                    PALETTE["slate"][600],  # Neutral gray
                    PALETTE["emerald"][600],  # Success green
                    PALETTE["amber"][600],  # Warning amber
                    PALETTE["indigo"][600],  # Deep blue
                    PALETTE["zinc"][500],  # Medium gray
                    PALETTE["sky"][600],  # Light blue
                    PALETTE["teal"][600],  # Teal accent
                ],
                "description": "Professional blues and grays for business presentations",
            },
            "vibrant": {
                "colors": [
                    PALETTE["blue"][500],  # Bright blue
                    PALETTE["pink"][500],  # Hot pink
                    PALETTE["green"][500],  # Vivid green
                    PALETTE["orange"][500],  # Bright orange
                    PALETTE["purple"][500],  # Rich purple
                    PALETTE["cyan"][500],  # Electric cyan
                    PALETTE["rose"][500],  # Rose red
                    PALETTE["lime"][500],  # Lime green
                ],
                "description": "Bold, energetic colors for creative presentations",
            },
            "pastel": {
                "colors": [
                    PALETTE["blue"][200],  # Soft blue
                    PALETTE["pink"][200],  # Soft pink
                    PALETTE["green"][200],  # Soft green
                    PALETTE["purple"][200],  # Soft purple
                    PALETTE["amber"][200],  # Soft amber
                    PALETTE["cyan"][200],  # Soft cyan
                    PALETTE["rose"][200],  # Soft rose
                    PALETTE["lime"][200],  # Soft lime
                ],
                "description": "Soft, muted colors for gentle emphasis",
            },
            "monochrome_blue": {
                "colors": [
                    PALETTE["blue"][900],  # Darkest
                    PALETTE["blue"][700],
                    PALETTE["blue"][600],
                    PALETTE["blue"][500],
                    PALETTE["blue"][400],
                    PALETTE["blue"][300],
                    PALETTE["blue"][200],  # Lightest
                ],
                "description": "Shades of blue for focused data visualization",
            },
            "monochrome_green": {
                "colors": [
                    PALETTE["green"][900],  # Darkest
                    PALETTE["green"][700],
                    PALETTE["green"][600],
                    PALETTE["green"][500],
                    PALETTE["green"][400],
                    PALETTE["green"][300],
                    PALETTE["green"][200],  # Lightest
                ],
                "description": "Shades of green for growth and positive metrics",
            },
            "earthy": {
                "colors": [
                    PALETTE["amber"][700],  # Deep amber/brown
                    PALETTE["orange"][600],  # Warm orange
                    PALETTE["lime"][700],  # Olive green
                    PALETTE["emerald"][700],  # Forest green
                    PALETTE["amber"][500],  # Golden
                    PALETTE["orange"][800],  # Dark orange
                    PALETTE["green"][800],  # Dark green
                ],
                "description": "Natural, warm tones for organic data",
            },
            "cool": {
                "colors": [
                    PALETTE["blue"][600],  # Cool blue
                    PALETTE["cyan"][600],  # Cyan
                    PALETTE["teal"][600],  # Teal
                    PALETTE["sky"][600],  # Sky blue
                    PALETTE["indigo"][600],  # Indigo
                    PALETTE["blue"][400],  # Light blue
                    PALETTE["cyan"][400],  # Light cyan
                ],
                "description": "Blues, cyans, and teals for calm, professional data",
            },
            "warm": {
                "colors": [
                    PALETTE["red"][600],  # Red
                    PALETTE["orange"][600],  # Orange
                    PALETTE["amber"][600],  # Amber
                    PALETTE["yellow"][600],  # Yellow
                    PALETTE["rose"][600],  # Rose
                    PALETTE["pink"][600],  # Pink
                    PALETTE["orange"][400],  # Light orange
                ],
                "description": "Reds, oranges, and yellows for energetic data",
            },
            "rainbow": {
                "colors": [
                    PALETTE["red"][500],  # Red
                    PALETTE["orange"][500],  # Orange
                    PALETTE["yellow"][500],  # Yellow
                    PALETTE["green"][500],  # Green
                    PALETTE["blue"][500],  # Blue
                    PALETTE["indigo"][500],  # Indigo
                    PALETTE["violet"][500],  # Violet
                    PALETTE["pink"][500],  # Pink
                ],
                "description": "Full spectrum for diverse categories",
            },
            "status": {
                "colors": [
                    PALETTE["green"][600],  # Success/Good
                    PALETTE["amber"][600],  # Warning/Caution
                    PALETTE["red"][600],  # Error/Critical
                    PALETTE["blue"][600],  # Info/Neutral
                    PALETTE["emerald"][500],  # Positive
                    PALETTE["orange"][600],  # Attention
                ],
                "description": "Traffic light colors for KPIs and status indicators",
            },
        }

        if style_preset not in style_presets:
            available = ", ".join(style_presets.keys())
            return json.dumps(
                {"error": f"Invalid style_preset '{style_preset}'. Available: {available}"},
                indent=2,
            )

        preset = style_presets[style_preset]

        return json.dumps(
            {
                "colors": preset["colors"],
                "description": preset["description"],
                "preset_name": style_preset,
                "usage": "Pass this to pptx_add_chart's options parameter: options={'colors': <colors array>}",
            },
            indent=2,
        )

    # Return the tools for external access
    return {"pptx_add_chart": pptx_add_chart, "pptx_get_chart_style": pptx_get_chart_style}
